#!/usr/bin/env python3
"""
Generate the Japanese manuscript for:
「疾患分類学的相対性：医療版サピア＝ウォーフ仮説の定式化」

Outputs:
  - manuscript_ja.docx  (日本語原稿・図表インライン)
  - figures_ja.pptx     (日本語図表・編集可能)
  - tables_ja.docx      (日本語表・編集可能)
"""

import os, re, textwrap
from pathlib import Path

from docx import Document
from docx.shared import Pt, Inches, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm

OUT_DIR = Path(__file__).resolve().parent.parent / "output"
OUT_DIR.mkdir(exist_ok=True)

# ---------------------------------------------------------------------------
# Japanese font setup for matplotlib
# ---------------------------------------------------------------------------

def setup_japanese_font():
    """Find and set a Japanese font for matplotlib."""
    jp_fonts = ['IPAGothic', 'IPAPGothic', 'Noto Sans CJK JP', 'TakaoPGothic',
                'VL PGothic', 'Takao PGothic', 'DejaVu Sans']
    for font_name in jp_fonts:
        try:
            fp = fm.findfont(fm.FontProperties(family=font_name), fallback_to_default=False)
            if fp and 'DejaVu' not in fp:
                plt.rcParams['font.family'] = font_name
                return font_name
        except Exception:
            continue
    # Fallback: use any available CJK font
    for f in fm.fontManager.ttflist:
        if any(kw in f.name.lower() for kw in ['gothic', 'mincho', 'cjk', 'ipa', 'noto']):
            plt.rcParams['font.family'] = f.name
            return f.name
    return None


# ---------------------------------------------------------------------------
# Utility helpers
# ---------------------------------------------------------------------------

def add_superscript_text(paragraph, text):
    parts = re.split(r'(\{[^}]+\})', text)
    for part in parts:
        if part.startswith('{') and part.endswith('}'):
            run = paragraph.add_run(part[1:-1])
            run.font.superscript = True
            run.font.size = Pt(8)
        else:
            run = paragraph.add_run(part)
            run.font.size = Pt(11)
    return paragraph


def set_cell_text(cell, text, bold=False, size=Pt(10)):
    cell.text = ""
    p = cell.paragraphs[0]
    run = p.add_run(text)
    run.font.size = size
    run.bold = bold


def add_heading(doc, text, level):
    h = doc.add_heading(text, level=level)
    for run in h.runs:
        run.font.color.rgb = RGBColor(0, 0, 0)
    return h


def add_body(doc, text):
    p = doc.add_paragraph()
    p.style = doc.styles['Normal']
    add_superscript_text(p, text)
    p.paragraph_format.space_after = Pt(6)
    p.paragraph_format.line_spacing = 1.5
    return p


# ---------------------------------------------------------------------------
# Figure generation (Japanese)
# ---------------------------------------------------------------------------

def create_figure1_ja():
    font_name = setup_japanese_font()
    fig, ax = plt.subplots(figsize=(10, 7))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 8)
    ax.axis('off')
    fig.patch.set_facecolor('white')

    ax.text(5, 7.6, '疾患分類学的相対性フレームワーク',
            ha='center', va='center', fontsize=14, fontweight='bold')

    box1 = dict(boxstyle='round,pad=0.4', facecolor='#E8F4FD', edgecolor='#2196F3', linewidth=2)
    ax.text(1.8, 6.2, 'レベル1：認知的\n（臨床推論）', ha='center', va='center',
            fontsize=10, bbox=box1)

    box2 = dict(boxstyle='round,pad=0.4', facecolor='#FFF3E0', edgecolor='#FF9800', linewidth=2)
    ax.text(5, 6.2, 'レベル2：制度的\n（研究・ガイドライン）', ha='center', va='center',
            fontsize=10, bbox=box2)

    box3 = dict(boxstyle='round,pad=0.4', facecolor='#E8F5E9', edgecolor='#4CAF50', linewidth=2)
    ax.text(8.2, 6.2, 'レベル3：集団的\n（疫学パターン）', ha='center', va='center',
            fontsize=10, bbox=box3)

    central = dict(boxstyle='round,pad=0.6', facecolor='#F3E5F5', edgecolor='#9C27B0', linewidth=2.5)
    ax.text(5, 4.0, '疾患分類学的\nフレームワーク\n（疾患分類体系）', ha='center', va='center',
            fontsize=11, fontweight='bold', bbox=central)

    arrow_style = dict(arrowstyle='->', color='#666666', lw=1.5,
                       connectionstyle='arc3,rad=0.1')
    ax.annotate('', xy=(1.8, 5.5), xytext=(4.0, 4.7), arrowprops=arrow_style)
    ax.annotate('', xy=(5, 5.5), xytext=(5, 4.7), arrowprops=arrow_style)
    ax.annotate('', xy=(8.2, 5.5), xytext=(6.0, 4.7), arrowprops=arrow_style)

    loop_style = dict(arrowstyle='->', color='#E91E63', lw=1.8,
                      connectionstyle='arc3,rad=-0.4', linestyle='dashed')
    ax.annotate('', xy=(3.8, 3.5), xytext=(1.5, 5.5), arrowprops=loop_style)
    ax.text(0.8, 4.5, 'ループ\n効果', ha='center', va='center',
            fontsize=8, color='#E91E63', fontstyle='italic')

    weak_box = dict(boxstyle='round,pad=0.3', facecolor='#FAFAFA', edgecolor='#999999', linewidth=1)
    ax.text(2.5, 1.8, '弱い形式\n（影響）', ha='center', va='center', fontsize=9, bbox=weak_box)
    ax.text(7.5, 1.8, '強い形式\n（決定）', ha='center', va='center', fontsize=9, bbox=weak_box)

    ax.annotate('', xy=(7.5, 1.2), xytext=(2.5, 1.2),
                arrowprops=dict(arrowstyle='<->', color='#333333', lw=1.5))
    ax.text(5, 0.9, '疾患分類学的制約のスペクトル', ha='center', va='center',
            fontsize=9, fontstyle='italic', color='#555555')

    counter_box = dict(boxstyle='round,pad=0.3', facecolor='#FFEBEE', edgecolor='#F44336', linewidth=1.5)
    ax.text(5, 2.8, '反証：分類非依存性疾患\n（例：痩身理想のない文化での拒食症）',
            ha='center', va='center', fontsize=8, bbox=counter_box)

    ax.annotate('', xy=(5, 3.3), xytext=(5, 3.0),
                arrowprops=dict(arrowstyle='->', color='#F44336', lw=1.2))

    fig_path = OUT_DIR / "figure1_framework_ja.png"
    fig.savefig(fig_path, dpi=300, bbox_inches='tight')
    plt.close(fig)
    return fig_path


def create_figure2_ja():
    font_name = setup_japanese_font()
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 5)
    ax.axis('off')
    fig.patch.set_facecolor('white')

    ax.text(5, 4.7, '疾患分類学的ループメカニズム', ha='center', fontsize=13, fontweight='bold')

    positions = [(2, 3.2), (8, 3.2), (8, 1.2), (2, 1.2)]
    labels = [
        '疾患分類学的\n分類体系\n（ICD, DSM等）',
        '臨床実践\n（診断・治療・\n紹介パターン）',
        '患者体験\n（自己同一化・\n症状表現）',
        '疫学的データ\nとエビデンス\n（有病率・転帰）'
    ]
    colors = ['#E8F4FD', '#FFF3E0', '#E8F5E9', '#FCE4EC']
    edge_colors = ['#2196F3', '#FF9800', '#4CAF50', '#E91E63']

    for (x, y), label, fc, ec in zip(positions, labels, colors, edge_colors):
        bbox = dict(boxstyle='round,pad=0.4', facecolor=fc, edgecolor=ec, linewidth=2)
        ax.text(x, y, label, ha='center', va='center', fontsize=9, bbox=bbox)

    arrow_kw = dict(arrowstyle='->', lw=2, color='#555555')
    ax.annotate('', xy=(6.3, 3.2), xytext=(3.7, 3.2), arrowprops=arrow_kw)
    ax.annotate('', xy=(8, 2.0), xytext=(8, 2.6), arrowprops=arrow_kw)
    ax.annotate('', xy=(3.7, 1.2), xytext=(6.3, 1.2), arrowprops=arrow_kw)
    ax.annotate('', xy=(2, 2.6), xytext=(2, 2.0), arrowprops=arrow_kw)

    ax.text(5, 3.5, '形成する', ha='center', fontsize=8, fontstyle='italic', color='#666')
    ax.text(8.6, 2.3, '規定する', ha='center', fontsize=8, fontstyle='italic', color='#666', rotation=90)
    ax.text(5, 0.9, '生成する', ha='center', fontsize=8, fontstyle='italic', color='#666')
    ax.text(1.4, 2.3, '改訂する', ha='center', fontsize=8, fontstyle='italic', color='#666', rotation=90)

    fig_path = OUT_DIR / "figure2_looping_ja.png"
    fig.savefig(fig_path, dpi=300, bbox_inches='tight')
    plt.close(fig)
    return fig_path


# ---------------------------------------------------------------------------
# Content
# ---------------------------------------------------------------------------

TITLE = "疾患分類学的相対性：医療版サピア＝ウォーフ仮説の定式化"
AUTHORS = "大西 達季"
AFFILIATIONS = ""

ABSTRACT = (
    "【背景】言語学のサピア＝ウォーフ仮説は、言語の構造が認知と知覚に影響を与えると主張する。"
    "医療においても類似の現象が存在する可能性がある：疾患分類学的フレームワーク（疾患の分類"
    "・命名体系）は、臨床的現実を単に記述するだけでなく、能動的にそれを形成している可能性がある。"
    "疾患概念が診療体系を規定するという主張は直感的に明らかに思えるが、その効果は自明の域を超えて"
    "おり、体系的な理論的取り扱いが必要であると我々は主張する。"
    "【目的】本論文は「疾患分類学的相対性」フレームワークを提案する。これは医療分類体系に適用された"
    "サピア＝ウォーフ仮説の形式的な類似物である。仮説の強い形式と弱い形式を定義し、3つのレベルの"
    "効果（認知的・制度的・集団的）を画定し、Ian Hackingのループ効果メカニズムを中核的なフィード"
    "バック動態として組み込む。"
    "【方法】科学哲学、医療人類学、診断の社会学、臨床医学からの知見を統合して理論的フレームワークを"
    "構築する。疼痛医学、産業保健（過労死）、ICD改定史からの症例研究を例証的根拠として分析する。"
    "分類非依存性の疾患発現（例：痩身理想のない社会における拒食症）からの反証を統合し、フレーム"
    "ワークの境界条件を画定する。"
    "【結果】疾患分類学的相対性フレームワークは6つの命題と3つの検証可能な予測を同定する。疾患分類"
    "学的カテゴリが臨床推論に対する認知的制約（診断的閉鎖）、研究資金配分と治療経路の制度的決定因、"
    "疫学パターンの集団レベルでの形成因として機能することを示す。"
    "【結論】疾患分類学的相対性フレームワークは、疾患分類体系がいかに診療を形成するかを理解するための"
    "構造化された理論的レンズを提供する。このフレームワークはICD改定前後の自然実験デザインに適した"
    "検証可能な予測を生成し、疾患分類学的政策と臨床教育に実践的含意を持つ。"
)

KEYWORDS = [
    "サピア＝ウォーフ仮説",
    "疾患分類学",
    "疾患分類",
    "言語的相対性",
    "ループ効果",
    "医学の哲学",
    "ICD",
    "診断ラベル",
]

REFERENCES = [
    "Sapir E. The status of linguistics as a science. Language. 1929;5(4):207–14.",
    "Whorf BL. Language, thought, and reality: selected writings of Benjamin Lee Whorf. Carroll JB, editor. Cambridge (MA): MIT Press; 1956.",
    "Warner R. The relationship between language and disease concepts. Int J Psychiatry Med. 1976;7(1):57–68.",
    "Hacking I. The looping effects of human kinds. In: Sperber D, Premack D, Premack AJ, editors. Causal cognition: a multidisciplinary debate. Oxford: Clarendon Press; 1995. p. 351–94.",
    "Hacking I. Making up people. London Review of Books. 2006;28(16):23–6.",
    "Kleinman A. The illness narratives: suffering, healing, and the human condition. New York: Basic Books; 1988.（邦訳：クラインマン A. 病いの語り. 誠信書房; 1996.）",
    "Jutel A. Sociology of diagnosis: a preliminary review. Sociol Health Illn. 2009;31(2):278–99.",
    "Zachar P, Kendler KS. The philosophy of nosology. Annu Rev Clin Psychol. 2017;13:49–71.",
    "Boorse C. Health as a theoretical concept. Philos Sci. 1977;44(4):542–73.",
    "Swartz L. Anorexia nervosa as a culture-bound syndrome. Soc Sci Med. 1985;20(7):725–30.",
    "Michaleff ZA, Glasziou P, Thomas R. Consequences of a diagnostic label: a systematic scoping review and thematic framework. Front Public Health. 2021;9:725877.",
    "Nickel B, Moynihan R, Barratt A, Brito JP, McCaffery K. Words do matter: a systematic review on how different terminology for the same condition influences management preferences. BMJ Open. 2017;7(7):e014129.",
    "Iwasaki K, Takahashi M, Nakata A. Health problems due to long working hours in Japan: working hours, workers' compensation (karoshi), and preventive measures. Ind Health. 2006;44(4):537–40.",
    "Nishiyama K, Johnson JV. Karoshi—death from overwork: occupational health consequences of Japanese production management. Int J Health Serv. 1997;27(4):625–41.",
    "World Health Organization. ICD-11 for Mortality and Morbidity Statistics. Geneva: WHO; 2019.",
    "Reed GM, First MB, Kogan CS, et al. Innovations and changes in the ICD-11 classification of mental, behavioural and neurodevelopmental disorders. World Psychiatry. 2019;18(1):3–19.",
    "Treede RD, Rief W, Barke A, et al. Chronic pain as a symptom or a disease: the IASP Classification of Chronic Pain for the International Classification of Diseases (ICD-11). Pain. 2019;160(1):19–27.",
    "Tsou JY. Natural kinds, psychiatric classification and the history of the DSM. Hist Psychiatry. 2016;27(4):406–24.",
    "Cooper R. Classifying madness: a philosophical examination of the Diagnostic and Statistical Manual of Mental Disorders. Dordrecht: Springer; 2005.",
    "Fabrega H Jr. Disease and social behavior: an interdisciplinary perspective. Cambridge (MA): MIT Press; 1974.",
    "Eisenberg L. Disease and illness: distinctions between professional and popular ideas of sickness. Cult Med Psychiatry. 1977;1(1):9–23.",
    "Rosenhan DL. On being sane in insane places. Science. 1973;179(4070):250–8.",
    "Thibault JM, Bhatt DL, Engel GL. The biopsychosocial model: past, present, future. Psychosomatics. 2003;44(4):267–75.",
    "Boroditsky L. Does language shape thought? Mandarin and English speakers' conceptions of time. Cogn Psychol. 2001;43(1):1–22.",
    "Lupyan G, Bergen B. How language programs the mind. Top Cogn Sci. 2016;8(2):408–24.",
    "Conrad P. The medicalization of society: on the transformation of human conditions into treatable disorders. Baltimore: Johns Hopkins University Press; 2007.",
    "Bowker GC, Star SL. Sorting things out: classification and its consequences. Cambridge (MA): MIT Press; 1999.",
    "Craddock N, Owen MJ. The Kraepelinian dichotomy — going, going... but still not gone. Br J Psychiatry. 2010;196(2):92–5.",
    "Frances A. Saving normal: an insider's revolt against out-of-control psychiatric diagnosis, DSM-5, Big Pharma, and the medicalization of ordinary life. New York: William Morrow; 2013.",
    "Kirmayer LJ. Cultural variations in the clinical presentation of depression and anxiety: implications for diagnosis and treatment. J Clin Psychiatry. 2001;62 Suppl 13:22–8.",
]


INTRO_PARAS = [
    (
        "サピア＝ウォーフ仮説は、20世紀初頭にEdward SapirとBenjamin Lee Whorfによって定式化"
        "された仮説であり、言語の構造がその話者の認知と世界観に影響を与えると主張する。{1,2}"
        "その強い形式（言語決定論）では言語が思考を決定し、弱い形式（言語相対性）では言語が"
        "習慣的な思考パターンに影響を与えるが完全には制約しないとする。{24,25}強い形式は"
        "言語学では大部分が放棄されているが、色彩知覚・空間認知・時間的認知を含む複数の"
        "領域において弱い形式を支持する堅固な証拠が存在する。{24}"
    ),
    (
        "医療においても類似の現象が作用している可能性がある。疾患分類学的フレームワーク"
        "——疾患が分類・命名・体系化される公式体系——は、臨床医学の「言語」を構成する。"
        "自然言語のカテゴリが話者の知覚と推論を形成するように、疾患分類学的カテゴリは"
        "臨床医が患者をいかに知覚し、診断し、治療するかを形成する可能性がある。"
        "Warnerは1976年にこの類推を初めて探求し、異なる文化の言語構造が根本的に異なる"
        "疾患概念につながると論じた。{3}しかし、この並行関係は初期の観察を超えて形式的に"
        "発展されていない。"
    ),
    (
        "「疾患概念が診療を形成する」という主張は自明に真であるように見えるかもしれない"
        "——医療は結局のところ疾患カテゴリを中心に組織されている。しかし我々は、その効果は"
        "自明の域をはるかに超えており、体系的な理論的取り扱いが必要であると主張する。"
        "この現象には、臨床推論に対する認知的効果（診断的閉鎖）、研究資金配分と治療"
        "ガイドラインに対する制度的効果、疫学パターンに対する集団レベルの効果が含まれる。"
        "さらに、Ian Hackingが精神医学的分類について示したように、特有のループメカニズムが"
        "作用する：分類体系が患者の自己同一化と症状表現を変化させ、それが翻って元の分類を"
        "補強するデータを生成する。{4,5}"
    ),
    (
        "同時に、この仮説を制約する証拠も存在する。Kleinmanは、痩せていることが理想とされ"
        "ない文化においても拒食症が発生することを観察しており、{6,10}一部の疾患は疾患分類"
        "学的枠組みとは独立に発現する生物学的基盤を有することを示唆する。適切な理論的フレーム"
        "ワークは、疾患分類学的カテゴリの構成的力と分類非依存性の疾患現象の存在の両方を"
        "包含しなければならない。"
    ),
    (
        "本論文では、「疾患分類学的相対性」（Nosological Relativity, NR）フレームワークを"
        "提案する——医療分類体系に適用されたサピア＝ウォーフ仮説の形式的類似物である。"
        "フレームワークの中核命題を定義し、強い形式と弱い形式を区別し、3つのレベルの効果を"
        "画定し、ループメカニズムを中核的フィードバック動態として定式化する。次に症例研究で"
        "フレームワークを例証し、実証的研究に適した検証可能な予測を導出する。"
    ),
]

BACKGROUND_SECTIONS = {
    "言語学におけるサピア＝ウォーフ仮説": [
        (
            "言語的相対性仮説は、その最初の定式化以来、相当な修正を経てきた。{1,2}現代の研究は"
            "中庸な立場を支持する：言語は思考を厳格に決定しないが、習慣的な認知パターンに影響を"
            "与える。{24,25}Boroditskyは、中国語と英語の話者が時間を異なって概念化することを"
            "示し、各言語が時間関係を符号化する構造的差異と一致する。{24}LupyanとBergenは、"
            "言語的ラベルがカテゴリ化を促進し知覚処理を調節することを示した。{25}これらの知見は、"
            "記号的分類体系が測定可能な認知的効果を発揮することを確立する——我々はこの原理を"
            "医療疾患分類学に拡張する。"
        ),
    ],
    "言語と疾患に関する先行研究": [
        (
            "Warnerの1976年の論文は、サピア＝ウォーフ仮説を医療に明示的に適用した最初の"
            "研究を代表する。{3}彼はインド＝ヨーロッパ語族の言語構造——特に疾病を記述する際の"
            "動詞よりも名詞の使用、空間メタファーの多用、主語＝述語の二分法——が静的、単因的、"
            "二元論的な疾患概念を促進すると論じた。Warnerは、これらの言語的特徴が外科手術への"
            "過度な依存を促進し、疾患原因における社会的・心理的因子の認識を阻害すると示唆した。"
        ),
        (
            "その後の研究は、定式化なしにこの観点を拡張してきた。Eisenbergは「疾患」（生物医学的"
            "構成概念）と「病い」（患者の生きた経験）を区別し、これらの構成概念間のギャップが"
            "文化的・言語的カテゴリによって媒介されることを示した。{21}Fabregaは疾患を"
            "社会的行動と結びつける学際的枠組みを提供した。{20}Kleinmanの医療人類学は、"
            "疾病の文化的カテゴリがラベリング、援助希求、治療的応答を導くことを実証した。{6}"
            "Hackingの人間種の哲学はループ効果概念を導入し、精神医学的分類が受動的なラベルでは"
            "なく、それが記述する現象の能動的構成要素であることを示した。{4,5}Jutelの診断の"
            "社会学は、診断行為自体が臨床的出会いを超えた帰結を持つ社会的事象であることを"
            "示した。{7}"
        ),
        (
            "これらの豊かな知的伝統にもかかわらず、疾患分類学的カテゴリが言語的相対性と類似した"
            "方法で診療を形成するという命題に対する統合的な形式的フレームワークは存在しない。"
            "本論文はこのギャップを埋める。"
        ),
    ],
}

FRAMEWORK_INTRO = (
    "ここで疾患分類学的相対性（NR）フレームワークを提示する。フレームワークは中核テーゼ、"
    "6つの形式的命題、3つのレベルの効果、ループメカニズムモデルから構成される（図1）。"
)

PROPOSITIONS = [
    (
        "命題1（中核テーゼ）。",
        "疾患分類学的フレームワーク——疾患が分類・命名される体系——は、既存の臨床的現実を"
        "単に記述するのではなく、医療システムの複数のレベルにおいて疾病の知覚・分類・管理を"
        "能動的に形成する。"
    ),
    (
        "命題2（弱い形式）。",
        "疾患分類学的カテゴリは臨床推論、研究優先事項、患者転帰に影響を与えるが、完全には"
        "決定しない。臨床医は訓練のカテゴリの外で推論することが可能であるが、そのような"
        "推論は追加の認知的努力を必要とし、統計的に頻度が低い。"
    ),
    (
        "命題3（強い形式）。",
        "特定の疾患——特に客観的バイオマーカーを欠くもの——については、疾患分類学的カテゴリが"
        "事実上臨床的現実を決定する：分類体系に名称が存在しない状態は診断・治療・研究の注目を"
        "受けず、したがって医療システム内で機能的に存在しないことになる。"
    ),
    (
        "命題4（認知的レベル）。",
        "疾患分類学的カテゴリは臨床医にとっての認知スキーマとして機能し、診断的閉鎖"
        "（カテゴリが付与された時点で診断推論が停止する現象）、カテゴリ整合的特徴への"
        "アンカリングバイアス、カテゴリ横断的パターンの知覚困難を生じさせる。"
    ),
    (
        "命題5（制度的レベル）。",
        "疾患分類学的カテゴリは研究資金配分、臨床ガイドライン策定、専門医研修、診療報酬"
        "体系を構造化する。疾患分類学的表現を持たない疾患は、その疾病負担に比して不釣り合いに"
        "少ない制度的支援を受ける。"
    ),
    (
        "命題6（集団レベル）。",
        "疾患分類学的カテゴリは、差別的な症例発見、診断的移行（カテゴリ変更時の既存症例の"
        "再分類）、援助希求行動を変化させる新たな疾患アイデンティティの創出を通じて、"
        "疫学パターンを形成する。"
    ),
]

LOOPING_PARAS = [
    (
        "NRフレームワークの中核はループメカニズムであり、Hackingの人間種の分析から適応した"
        "ものである（図2）。{4,5}メカニズムは以下のように作用する："
    ),
    (
        "段階1（分類→実践）：疾患分類学的カテゴリが公式体系（例：ICD, DSM）内で導入または"
        "改訂される。これにより、カテゴリを中心に組織された診断基準、治療アルゴリズム、"
        "紹介経路が創出される。"
    ),
    (
        "段階2（実践→体験）：カテゴリを使用する臨床医が患者に診断を伝達し、患者はそれを"
        "自己理解に組み込む。診断ラベルは症状解釈の枠組み、疾患アイデンティティの基盤、"
        "社会的役割を提供する。{11}"
    ),
    (
        "段階3（体験→データ）：診断カテゴリに自己同一化した患者は、カテゴリを使用して"
        "コード化される医療希求行動、症状報告、臨床的出会いを生成する。これによりカテゴリの"
        "存在と整合的な疫学的データが産出される。{26}"
    ),
    (
        "段階4（データ→分類）：疫学的データが疾患分類学的カテゴリの妥当性確認と精緻化に"
        "使用される。定義と整合的なデータを生成するカテゴリは経験的に支持されているように"
        "見え、ループが完結する。"
    ),
    (
        "このループメカニズムは自己強化的なサイクルを創出し、「自然の関節で切る」カテゴリと、"
        "それが記述すると主張するまさにそのパターンを生成するカテゴリとの区別を困難にする。{18,19}"
    ),
]

EVIDENCE_SECTIONS = {
    "疼痛医学における診断的閉鎖": [
        (
            "ICD-11における独立した疾患実体としての慢性疼痛の認知（コードMG30）は、疾患分類学的"
            "カテゴリの認知的レベルの効果を例証する。{17}ICD-11以前、慢性疼痛は基礎疾患の症状として"
            "もっぱら分類されていた。この疾患分類学的枠組みは臨床推論を制約した：基礎疾患が同定されると"
            "（例：変形性関節症）、疼痛メカニズムの診断的探索は典型的に停止した——我々がこれを"
            "診断的閉鎖と呼ぶ現象である。{4}"
        ),
        (
            "独立した診断実体としての慢性疼痛の導入は自然実験の条件を創出する。NRフレームワークは、"
            "ICD-11の慢性疼痛カテゴリの採用後、臨床医が：(a)基礎疾患のレンズのみを通じて疼痛を"
            "治療するのではなく、より多くの患者を疼痛専門医に紹介する；(b)疼痛特異的治療の処方を"
            "増加させる；(c)疾患分類学的変更ではなく検出の改善に帰される慢性疼痛有病率に関する"
            "新たな疫学的データを生成すると予測する。"
        ),
        (
            "Nickelらはシステマティックレビューにおいて、同一疾患に対する異なる用語が臨床医と"
            "患者の管理選好に影響を与えることを示した。{12}この知見は命題2を直接支持する："
            "疾患分類学的ラベルは臨床推論に影響を与えるが、決定はしない。"
        ),
    ],
    "過労死：文化特異的な疾患分類学的革新": [
        (
            "過労死は命題3——疾患分類学的相対性の強い形式——の説得的な例証を提供する。この概念は"
            "1970年代の日本で、基礎にある心血管または脳血管イベントとは区別された、法的・医学的に"
            "認知された死因として出現した。{13,14}他の国々では、過重労働の文脈で発生する同一の"
            "病態生理学的イベント（心筋梗塞、脳卒中）は、職業的因果関係への言及なしに臓器特異的"
            "カテゴリの下に分類される。"
        ),
        (
            "日本における疾患分類学的カテゴリとしての過労死の存在は、NRフレームワークと整合的な"
            "効果のカスケードを生じさせてきた：(a)認知的レベルでは、日本の臨床医は労働年齢成人の"
            "突然の心血管死の鑑別診断において職業的過重労働を考慮する——等価のカテゴリを持たない"
            "国々では大部分不在の考慮である；(b)制度的レベルでは、日本は過労死認定のための専用の"
            "法医学的装置を発展させ、他国には存在しない特定の残業時間閾値（「80時間の過労死ライン」）"
            "を設定している；(c)集団レベルでは、日本は過労死統計を独自の疫学的現象として報告し、"
            "カテゴリが存在しない国々では構造的に不可能な政策的応答（2014年過労死等防止対策"
            "推進法）を生成している。"
        ),
        (
            "ICD-11における分類可能な職業的現象としてのバーンアウト（QD85）の収載は、この"
            "疾患分類学的革新の部分的な国際的拡散を提供し、NRフレームワークはこのコードを"
            "採用する国々が——より小さな規模で——日本で観察された制度的・疫学的パターンを"
            "再現し始めると予測する。{15,16}"
        ),
    ],
    "自然実験としてのICD改定": [
        (
            "国際疾病分類の大規模改定はNRフレームワークの自然実験を構成する。ICD-10からICD-11への"
            "移行はいくつかの疾患分類学的革新を導入し、その効果は前後比較デザインを用いて追跡"
            "可能である（表1）。{15,16}"
        ),
        (
            "理論的に特に興味深いのは、性別不合の精神および行動の障害（ICD-10, 第V章）から"
            "性の健康に関連する状態（ICD-11, 第17章）への再分類である。NRフレームワークは、"
            "この再分類が3つのレベルすべてで測定可能な変化を生じさせると予測する：臨床医における"
            "精神疾患との認知的連合の減少、精神科サービスから離れたケア経路の制度的再編成、"
            "スティグマ関連のケアアクセス障壁の集団レベルでの減少。{16}"
        ),
    ],
    "反証：分類非依存性疾患": [
        (
            "NRフレームワークは、一部の疾患が疾患分類学的枠組みとは独立して発現するという"
            "証拠を包含しなければならない。Kleinmanは痩せていることが理想とされない文化において"
            "拒食症が発生することを記録し、{6}Swartzは拒食症の中核的特徴が持続する場合でも"
            "その形態が文化間で変動するがゆえに、拒食症は文化結合症候群として理解されるべきだと"
            "論じた。{10}Kirmayerはさらに、うつ病や不安の臨床的発現における文化的変異が"
            "単純な普遍主義的仮定を複雑にすることを記録している。{30}より広くは、明確な"
            "病因因子を持つ感染症（例：結核、マラリア）は、それがどのように分類されるかに"
            "かかわらず発現する。"
        ),
        (
            "これらの観察は疾患分類学的相対性の境界条件を画定する。我々は、疾患分類学的効果の"
            "強度は基礎にある生物学的シグナルの強度と反比例すると提案する：強い、客観的に測定"
            "可能な生物学的基盤を持つ疾患（例：骨折、同定可能な病原体を持つ感染症）は疾患分類学的"
            "形成に対する感受性が低い一方、主に症状群、機能障害、主観的体験によって定義される疾患"
            "（例：線維筋痛症、慢性疲労症候群、多くの精神障害）は高い感受性を持つ。{9}この勾配は"
            "NRフレームワークの検証可能な予測を構成する。"
        ),
    ],
}

TABLE1_CAPTION = "表1. ICD-11における主要な疾患分類学的革新と疾患分類学的相対性フレームワークにおける予測効果"
TABLE1_HEADERS = ["ICD-11の変更", "認知的レベル", "制度的レベル", "集団レベル"]
TABLE1_ROWS = [
    [
        "独立した疾患実体としての慢性疼痛（MG30）",
        "臨床医が疼痛を症状ではなく一次的診断として考慮",
        "疼痛医学が専用の資金源とガイドラインを獲得",
        "検出の改善に帰される慢性疼痛有病率の増加",
    ],
    [
        "職業的現象としてのバーンアウト（QD85）",
        "臨床医が疲弊の症例で職業的因子をスクリーニング",
        "産業保健サービスがバーンアウト特異的プロトコルを策定",
        "バーンアウト有病率に関する新たな疫学データの出現",
    ],
    [
        "ゲーム障害（6C51）",
        "臨床医が病的ゲームを診断可能な状態として同定",
        "ゲーム障害の治療センターと臨床ガイドラインの発展",
        "ゲーム障害の発生率データの生成と報告",
    ],
    [
        "複雑性PTSD（6B41）",
        "臨床医がトラウマ患者で複雑性と単純性のPTSDを区別",
        "複雑性PTSDメカニズムに特化した研究資金の配分",
        "既存PTSD患者の再分類による有病率推定値の変化",
    ],
    [
        "性別不合の再分類（HA60–HA6Z）",
        "精神疾患との認知的連合の減少",
        "精神科サービスから離れたケア経路の再編成",
        "スティグマ関連のケアアクセス障壁の減少",
    ],
]

PREDICTIONS_PARAS = [
    (
        "NRフレームワークは3つの経験的に検証可能な予測を生成する："
    ),
    (
        "予測1（ICD改定効果）：ICD-11における新たな診断カテゴリの導入は、基礎にある"
        "生物学的発生率のいかなる変化とも独立に、対応する疾患の検出率、専門医紹介、"
        "標的治療処方の増加を伴う。これはICD移行日前後のレセプトデータの中断時系列分析を"
        "用いて検証可能である。"
    ),
    (
        "予測2（バイオマーカー勾配）：診断・治療パターンに対する疾患分類学的効果の規模は、"
        "当該疾患に対する客観的バイオマーカーの利用可能性と特異度と逆相関する。もっぱら"
        "臨床基準で診断される疾患は、確定的な検査・画像マーカーを持つ疾患よりも大きな"
        "効果を示す。"
    ),
    (
        "予測3（国際間の乖離）：同一の病態生理学的現象に対して異なる疾患分類学的伝統を"
        "持つ国々（例：日本の過労死対他国の臓器特異的コーディング）は、それらの疾患に"
        "対する診断、治療、産業保健政策、患者の自己同一化において体系的に異なるパターンを"
        "示す。"
    ),
]

DISCUSSION_PARAS = [
    (
        "疾患分類学的相対性フレームワークは、診療と疾患分類学的政策にいくつかの含意を提供する。"
    ),
    (
        "第一に、疾患分類学的改定は、臨床的現実の改善された記述としてのみではなく、臨床実践を"
        "能動的に再形成する介入として理解されるべきである。これは、分類体系の改定プロセスに、"
        "法規制の規制影響評価に類似した前向きの影響評価を組み込むべきことを示唆する。{8,27}"
    ),
    (
        "第二に、臨床教育は診断カテゴリの認知的効果に関する明示的な教育を含むべきである。"
        "診断的閉鎖——カテゴリが付与された時点で診断推論を停止する傾向——の認知は、他の認知的"
        "ヒューリスティックのデバイアシング訓練と類似して、その効果を部分的に緩和しうる。{22}"
        "生物心理社会的アプローチとの統合は、カテゴリ的思考に内在する還元主義をさらに"
        "打ち消す可能性がある。{23}"
    ),
    (
        "第三に、フレームワークは新たな診断カテゴリを創出する提案を評価するための原則的基盤を"
        "提供する。疾患分類学的認知の利益（検出の改善、研究資金、治療開発）は、実体化の"
        "リスク（暫定的カテゴリを自然種として扱うこと）およびループ効果のリスク（カテゴリが"
        "それが記述するまさにその現象を生成すること）と衡量されなければならない。{26,29}"
    ),
    (
        "本フレームワークの限界を認めるべきである。NRフレームワークは現段階では主に理論的"
        "構成概念であり、上述の検証可能な予測は経験的妥当性確認を必要とする。自然言語と"
        "疾患分類学的体系の間の類推は不完全である：自然言語は暗黙的に習得され集団全体で"
        "使用されるが、疾患分類学的体系は明示的に構築され主に専門家によって使用される。"
        "さらに、フレームワークは疾患分類学的カテゴリの構成的効果を最小化すべきか活用すべきか"
        "という規範的問題には取り組まない。{8}"
    ),
]

CONCLUSION_PARAS = [
    (
        "我々は疾患分類体系がいかに診療を形成するかを理解するための形式的な理論的レンズとして"
        "疾患分類学的相対性フレームワークを提案した。言語学のサピア＝ウォーフ仮説、科学哲学の"
        "Hackingのループ効果メカニズム、医療人類学と診断の社会学からの証拠に基づき、6つの"
        "命題、3つのレベルの効果、3つの検証可能な予測を定義した。"
    ),
    (
        "フレームワークは生物学的疾患基盤の存在を否定しない；むしろ、これらの基盤が知覚・"
        "分類・管理される疾患分類学的「レンズ」が体系的で予測可能な歪みを導入すると論じる。"
        "これらの歪みを明示的かつ検証可能にすることにより、NRフレームワークは実証的研究の"
        "新たな方途を開き、疾患分類学的政策に原則的な指針を提供する。{27,28}"
    ),
    (
        "今後の研究は、上述の3つの予測の経験的検証——特にICD-10からICD-11への移行を自然"
        "実験として用いること——を優先すべきである。疾患分類学的伝統が乖離する疾患（過労死"
        "など）の国際比較研究は、追加の有望な研究デザインを提供する。究極的に、疾患分類学的"
        "相対性フレームワークは医療における反省的転回を促す：疾患を分類する我々の体系は"
        "生物学的現実の中立的な鏡ではなく、臨床知識の構築における能動的な参与者であるという"
        "認識である。"
    ),
]


# ---------------------------------------------------------------------------
# Build the Japanese manuscript
# ---------------------------------------------------------------------------

def build_manuscript(fig1_path, fig2_path):
    doc = Document()
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Times New Roman'
    font.size = Pt(11)
    style.paragraph_format.line_spacing = 1.5

    # Title page
    tp = doc.add_paragraph()
    tp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = tp.add_run(TITLE)
    run.bold = True
    run.font.size = Pt(16)

    tp2 = doc.add_paragraph()
    tp2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run2 = tp2.add_run(AUTHORS)
    run2.font.size = Pt(12)

    doc.add_paragraph()

    # Abstract
    add_heading(doc, '抄録', level=1)
    add_body(doc, ABSTRACT)
    kw_para = doc.add_paragraph()
    kw_run = kw_para.add_run('キーワード：')
    kw_run.bold = True
    kw_run.font.size = Pt(11)
    kw_text = kw_para.add_run('；'.join(KEYWORDS))
    kw_text.font.size = Pt(11)

    doc.add_page_break()

    # 1. Introduction
    add_heading(doc, '1. 緒言', level=1)
    for para_text in INTRO_PARAS:
        add_body(doc, para_text)

    # 2. Background
    add_heading(doc, '2. 背景', level=1)
    sec_num = 1
    for sub_title, paras in BACKGROUND_SECTIONS.items():
        add_heading(doc, f'2.{sec_num} {sub_title}', level=2)
        sec_num += 1
        for para_text in paras:
            add_body(doc, para_text)

    # 3. Framework
    add_heading(doc, '3. 疾患分類学的相対性フレームワーク', level=1)
    add_body(doc, FRAMEWORK_INTRO)

    # Figure 1
    doc.add_paragraph()
    fig1_para = doc.add_paragraph()
    fig1_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = fig1_para.add_run()
    run.add_picture(str(fig1_path), width=Inches(5.5))
    cap1 = doc.add_paragraph()
    cap1.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cap1.paragraph_format.space_before = Pt(12)
    cap_run = cap1.add_run(
        '図1. 疾患分類学的相対性フレームワーク。疾患分類学的カテゴリは'
        'ループメカニズムを通じて3つのレベル（認知的・制度的・集団的）で'
        '診療を形成する。効果の強度は弱い形式（影響）から強い形式（決定）'
        'までのスペクトルに沿って変動する。'
    )
    cap_run.italic = True
    cap_run.font.size = Pt(10)
    doc.add_paragraph()

    # Propositions
    add_heading(doc, '3.1 中核命題', level=2)
    for prop_title, prop_text in PROPOSITIONS:
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(6)
        p.paragraph_format.line_spacing = 1.5
        run_t = p.add_run(prop_title + ' ')
        run_t.bold = True
        run_t.font.size = Pt(11)
        run_b = p.add_run(prop_text)
        run_b.font.size = Pt(11)

    # Looping mechanism
    add_heading(doc, '3.2 ループメカニズム', level=2)
    for para_text in LOOPING_PARAS:
        add_body(doc, para_text)

    # Figure 2
    doc.add_paragraph()
    fig2_para = doc.add_paragraph()
    fig2_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = fig2_para.add_run()
    run.add_picture(str(fig2_path), width=Inches(5.5))
    cap2 = doc.add_paragraph()
    cap2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cap2.paragraph_format.space_before = Pt(12)
    cap_run2 = cap2.add_run(
        '図2. 疾患分類学的ループメカニズム。分類→実践→体験→データ→分類の'
        '4段階サイクル。この自己強化的ループが疾患分類学的相対性フレームワークの'
        '中核動態である。'
    )
    cap_run2.italic = True
    cap_run2.font.size = Pt(10)
    doc.add_paragraph()

    # 4. Evidence
    add_heading(doc, '4. 根拠と症例研究', level=1)
    sec_num = 1
    for sub_title, paras in EVIDENCE_SECTIONS.items():
        add_heading(doc, f'4.{sec_num} {sub_title}', level=2)
        sec_num += 1
        for i, para_text in enumerate(paras):
            add_body(doc, para_text)
            if '自然実験としてのICD改定' in sub_title and i == 0:
                doc.add_paragraph()
                t1cap = doc.add_paragraph()
                t1cap.paragraph_format.space_before = Pt(14)
                t1_run = t1cap.add_run(TABLE1_CAPTION)
                t1_run.bold = True
                t1_run.font.size = Pt(10)
                table = doc.add_table(rows=1 + len(TABLE1_ROWS), cols=4)
                table.style = 'Table Grid'
                table.alignment = WD_TABLE_ALIGNMENT.CENTER
                for j, header in enumerate(TABLE1_HEADERS):
                    set_cell_text(table.rows[0].cells[j], header, bold=True, size=Pt(9))
                for r_idx, row_data in enumerate(TABLE1_ROWS):
                    for c_idx, cell_text in enumerate(row_data):
                        set_cell_text(table.rows[r_idx + 1].cells[c_idx], cell_text, size=Pt(9))
                widths = [Inches(1.6), Inches(1.6), Inches(1.6), Inches(1.6)]
                for row in table.rows:
                    for idx, width in enumerate(widths):
                        row.cells[idx].width = width
                doc.add_paragraph()

    # 5. Predictions
    add_heading(doc, '5. 検証可能な予測', level=1)
    for para_text in PREDICTIONS_PARAS:
        add_body(doc, para_text)

    # 6. Discussion
    add_heading(doc, '6. 考察', level=1)
    for para_text in DISCUSSION_PARAS:
        add_body(doc, para_text)

    # 7. Conclusion
    add_heading(doc, '7. 結論', level=1)
    for para_text in CONCLUSION_PARAS:
        add_body(doc, para_text)

    # References
    doc.add_page_break()
    add_heading(doc, '参考文献', level=1)
    for i, ref in enumerate(REFERENCES, 1):
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(2)
        p.paragraph_format.line_spacing = 1.15
        run_num = p.add_run(f'{i}. ')
        run_num.bold = True
        run_num.font.size = Pt(10)
        run_ref = p.add_run(ref)
        run_ref.font.size = Pt(10)

    out_path = OUT_DIR / "manuscript_ja.docx"
    doc.save(str(out_path))
    print(f"日本語原稿保存: {out_path}")
    return out_path


def build_figures_pptx(fig1_path, fig2_path):
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    for fig_path, title_text, caption_text in [
        (fig1_path,
         "図1. 疾患分類学的相対性フレームワーク",
         "疾患分類学的カテゴリはループメカニズムを通じて3つのレベルで診療を形成する。"),
        (fig2_path,
         "図2. 疾患分類学的ループメカニズム",
         "分類→実践→体験→データ→分類の4段階サイクル。"),
    ]:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3),
                                          PptxInches(12.333), PptxInches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = PptxPt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        slide.shapes.add_picture(str(fig_path), PptxInches(1.5), PptxInches(1.3), width=PptxInches(10.333))

        capBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(6.3),
                                           PptxInches(12.333), PptxInches(1.0))
        cf = capBox.text_frame
        cf.word_wrap = True
        cp = cf.paragraphs[0]
        cp.text = caption_text
        cp.font.size = PptxPt(14)
        cp.font.italic = True
        cp.alignment = PP_ALIGN.CENTER

    out_path = OUT_DIR / "figures_ja.pptx"
    prs.save(str(out_path))
    print(f"日本語図表PPTX保存: {out_path}")
    return out_path


def build_tables_docx():
    doc = Document()
    style = doc.styles['Normal']
    style.font.name = 'Times New Roman'
    style.font.size = Pt(10)
    add_heading(doc, TABLE1_CAPTION, level=2)
    table = doc.add_table(rows=1 + len(TABLE1_ROWS), cols=4)
    table.style = 'Table Grid'
    for j, header in enumerate(TABLE1_HEADERS):
        set_cell_text(table.rows[0].cells[j], header, bold=True, size=Pt(10))
    for r_idx, row_data in enumerate(TABLE1_ROWS):
        for c_idx, cell_text in enumerate(row_data):
            set_cell_text(table.rows[r_idx + 1].cells[c_idx], cell_text, size=Pt(10))
    out_path = OUT_DIR / "tables_ja.docx"
    doc.save(str(out_path))
    print(f"日本語表DOCX保存: {out_path}")
    return out_path


def main():
    print("日本語図表を生成中...")
    fig1 = create_figure1_ja()
    fig2 = create_figure2_ja()
    print("日本語原稿を構築中...")
    build_manuscript(fig1, fig2)
    print("日本語図表PPTXを構築中...")
    build_figures_pptx(fig1, fig2)
    print("日本語表DOCXを構築中...")
    build_tables_docx()
    print("完了!")


if __name__ == '__main__':
    main()
