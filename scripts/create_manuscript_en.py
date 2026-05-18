#!/usr/bin/env python3
"""
Generate the English manuscript for:
"Nosological Relativity: Formalizing the Medical Sapir-Whorf Hypothesis"

Outputs:
  - manuscript_en.docx  (full manuscript with inline figures/tables)
  - figures_en.pptx     (editable figures, one per slide)
  - tables_en.docx      (editable tables)
"""

import os, re, textwrap
from pathlib import Path

from docx import Document
from docx.shared import Pt, Inches, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptxRGB

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch

OUT_DIR = Path(__file__).resolve().parent.parent / "output"
OUT_DIR.mkdir(exist_ok=True)

# ---------------------------------------------------------------------------
# Utility helpers
# ---------------------------------------------------------------------------

def add_superscript_text(paragraph, text):
    """Parse text with {n} or {n-m} markers and render them as Word-native superscripts."""
    parts = re.split(r'(\{[^}]+\})', text)
    for part in parts:
        if part.startswith('{') and part.endswith('}'):
            run = paragraph.add_run(part[1:-1])
            run.font.superscript = True
            run.font.size = Pt(8)
        else:
            run = paragraph.add_run(part)
            run.font.size = Pt(11)
    return paragraph


def set_cell_text(cell, text, bold=False, size=Pt(10)):
    cell.text = ""
    p = cell.paragraphs[0]
    run = p.add_run(text)
    run.font.size = size
    run.bold = bold


def add_heading(doc, text, level):
    h = doc.add_heading(text, level=level)
    for run in h.runs:
        run.font.color.rgb = RGBColor(0, 0, 0)
    return h


def add_body(doc, text):
    """Add a body paragraph with superscript citation support."""
    p = doc.add_paragraph()
    p.style = doc.styles['Normal']
    add_superscript_text(p, text)
    p.paragraph_format.space_after = Pt(6)
    p.paragraph_format.line_spacing = 1.5
    return p


# ---------------------------------------------------------------------------
# Figure generation (matplotlib)
# ---------------------------------------------------------------------------

def create_figure1():
    """Nosological Relativity Framework — conceptual diagram."""
    fig, ax = plt.subplots(figsize=(10, 7))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 8)
    ax.axis('off')
    fig.patch.set_facecolor('white')

    # Title
    ax.text(5, 7.6, 'Nosological Relativity Framework',
            ha='center', va='center', fontsize=14, fontweight='bold')

    # --- Three levels (boxes) ---
    box_props = dict(boxstyle='round,pad=0.4', facecolor='#E8F4FD', edgecolor='#2196F3', linewidth=2)
    ax.text(1.8, 6.2, 'Level 1: Cognitive\n(Clinician reasoning)', ha='center', va='center',
            fontsize=10, bbox=box_props)

    box_props2 = dict(boxstyle='round,pad=0.4', facecolor='#FFF3E0', edgecolor='#FF9800', linewidth=2)
    ax.text(5, 6.2, 'Level 2: Institutional\n(Research & guidelines)', ha='center', va='center',
            fontsize=10, bbox=box_props2)

    box_props3 = dict(boxstyle='round,pad=0.4', facecolor='#E8F5E9', edgecolor='#4CAF50', linewidth=2)
    ax.text(8.2, 6.2, 'Level 3: Population\n(Epidemiological patterns)', ha='center', va='center',
            fontsize=10, bbox=box_props3)

    # --- Central construct ---
    central = dict(boxstyle='round,pad=0.6', facecolor='#F3E5F5', edgecolor='#9C27B0', linewidth=2.5)
    ax.text(5, 4.0, 'NOSOLOGICAL\nFRAMEWORK\n(Disease Classification)', ha='center', va='center',
            fontsize=11, fontweight='bold', bbox=central)

    # --- Arrows from central to levels ---
    arrow_style = dict(arrowstyle='->', color='#666666', lw=1.5,
                       connectionstyle='arc3,rad=0.1')
    ax.annotate('', xy=(1.8, 5.5), xytext=(4.0, 4.7),
                arrowprops=arrow_style)
    ax.annotate('', xy=(5, 5.5), xytext=(5, 4.7),
                arrowprops=arrow_style)
    ax.annotate('', xy=(8.2, 5.5), xytext=(6.0, 4.7),
                arrowprops=arrow_style)

    # --- Looping effect (feedback) ---
    loop_style = dict(arrowstyle='->', color='#E91E63', lw=1.8,
                      connectionstyle='arc3,rad=-0.4', linestyle='dashed')
    ax.annotate('', xy=(3.8, 3.5), xytext=(1.5, 5.5),
                arrowprops=loop_style)
    ax.text(0.8, 4.5, 'Looping\neffect', ha='center', va='center',
            fontsize=8, color='#E91E63', fontstyle='italic')

    # --- Bottom: Strong vs Weak ---
    weak_box = dict(boxstyle='round,pad=0.3', facecolor='#FAFAFA', edgecolor='#999999', linewidth=1)
    ax.text(2.5, 1.8, 'Weak Form\n(Influence)', ha='center', va='center',
            fontsize=9, bbox=weak_box)
    ax.text(7.5, 1.8, 'Strong Form\n(Determination)', ha='center', va='center',
            fontsize=9, bbox=weak_box)

    # Spectrum arrow
    ax.annotate('', xy=(7.5, 1.2), xytext=(2.5, 1.2),
                arrowprops=dict(arrowstyle='<->', color='#333333', lw=1.5))
    ax.text(5, 0.9, 'Spectrum of nosological constraint', ha='center', va='center',
            fontsize=9, fontstyle='italic', color='#555555')

    # --- Counter-evidence box ---
    counter_box = dict(boxstyle='round,pad=0.3', facecolor='#FFEBEE', edgecolor='#F44336', linewidth=1.5)
    ax.text(5, 2.8, 'Counter-evidence:\nDisease-independent-of-nosology\n(e.g., anorexia in non-thin-ideal cultures)',
            ha='center', va='center', fontsize=8, bbox=counter_box)

    # Arrow from counter to central
    ax.annotate('', xy=(5, 3.3), xytext=(5, 3.0),
                arrowprops=dict(arrowstyle='->', color='#F44336', lw=1.2))

    fig_path = OUT_DIR / "figure1_framework.png"
    fig.savefig(fig_path, dpi=300, bbox_inches='tight')
    plt.close(fig)
    return fig_path


def create_figure2():
    """Looping effect mechanism diagram."""
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 5)
    ax.axis('off')
    fig.patch.set_facecolor('white')

    ax.text(5, 4.7, 'The Nosological Looping Mechanism', ha='center', fontsize=13, fontweight='bold')

    # Four boxes in a cycle
    positions = [(2, 3.2), (8, 3.2), (8, 1.2), (2, 1.2)]
    labels = [
        'Nosological\nClassification\n(e.g., ICD, DSM)',
        'Clinical Practice\n(Diagnosis, Treatment,\nReferral patterns)',
        'Patient Experience\n(Self-identification,\nSymptom expression)',
        'Epidemiological\nData & Evidence\n(Prevalence, Outcomes)'
    ]
    colors = ['#E8F4FD', '#FFF3E0', '#E8F5E9', '#FCE4EC']
    edge_colors = ['#2196F3', '#FF9800', '#4CAF50', '#E91E63']

    for (x, y), label, fc, ec in zip(positions, labels, colors, edge_colors):
        bbox = dict(boxstyle='round,pad=0.4', facecolor=fc, edgecolor=ec, linewidth=2)
        ax.text(x, y, label, ha='center', va='center', fontsize=9, bbox=bbox)

    # Arrows forming the cycle
    arrow_kw = dict(arrowstyle='->', lw=2, color='#555555')
    ax.annotate('', xy=(6.3, 3.2), xytext=(3.7, 3.2), arrowprops=arrow_kw)  # 1→2
    ax.annotate('', xy=(8, 2.0), xytext=(8, 2.6), arrowprops=arrow_kw)       # 2→3
    ax.annotate('', xy=(3.7, 1.2), xytext=(6.3, 1.2), arrowprops=arrow_kw)  # 3→4
    ax.annotate('', xy=(2, 2.6), xytext=(2, 2.0), arrowprops=arrow_kw)       # 4→1

    # Label the arrows
    ax.text(5, 3.5, 'Shapes', ha='center', fontsize=8, fontstyle='italic', color='#666')
    ax.text(8.6, 2.3, 'Molds', ha='center', fontsize=8, fontstyle='italic', color='#666', rotation=90)
    ax.text(5, 0.9, 'Generates', ha='center', fontsize=8, fontstyle='italic', color='#666')
    ax.text(1.4, 2.3, 'Revises', ha='center', fontsize=8, fontstyle='italic', color='#666', rotation=90)

    fig_path = OUT_DIR / "figure2_looping.png"
    fig.savefig(fig_path, dpi=300, bbox_inches='tight')
    plt.close(fig)
    return fig_path


# ---------------------------------------------------------------------------
# Manuscript content
# ---------------------------------------------------------------------------

TITLE = "Nosological Relativity: Formalizing the Medical Sapir–Whorf Hypothesis"

AUTHORS = "Tatsuki Onishi"
AFFILIATIONS = ""  # To be filled by the author

ABSTRACT = (
    "Background: The Sapir–Whorf hypothesis in linguistics posits that the structure of language "
    "influences cognition and perception. An analogous phenomenon may operate in medicine: "
    "nosological frameworks—the systems by which diseases are classified and named—may not merely "
    "describe clinical reality but actively shape it. While the notion that disease concepts "
    "influence medical practice is intuitive, we argue that the effects extend beyond the obvious "
    "and warrant formal theoretical treatment. "
    "Objective: This paper proposes the 'Nosological Relativity' framework, a formal analogue of "
    "the Sapir–Whorf hypothesis applied to medical classification systems. We define strong and "
    "weak forms of the hypothesis, delineate three levels of effect (cognitive, institutional, and "
    "population), and incorporate Ian Hacking's looping-effect mechanism as the central feedback "
    "dynamic. "
    "Methods: We synthesize insights from philosophy of science, medical anthropology, sociology "
    "of diagnosis, and clinical medicine to construct the theoretical framework. Case studies from "
    "pain medicine, occupational health (karoshi), and ICD revision history are analyzed as "
    "illustrative evidence. Counter-evidence from culture-independent disease manifestations "
    "(e.g., anorexia nervosa in non-thin-ideal societies) is integrated to delimit the "
    "framework's boundaries. "
    "Results: The Nosological Relativity framework identifies six propositions and three testable "
    "predictions. We demonstrate that nosological categories operate as cognitive constraints on "
    "clinical reasoning (diagnostic foreclosure), institutional determinants of research funding "
    "and therapeutic pathways, and population-level shapers of epidemiological patterns. The "
    "looping mechanism—whereby classification alters patient self-identification, which in turn "
    "generates data that reinforces the classification—is formalized as the central dynamic. "
    "Conclusions: The Nosological Relativity framework provides a structured theoretical lens for "
    "understanding how disease classification systems shape medical practice. The framework yields "
    "testable predictions amenable to natural-experiment designs, particularly around ICD "
    "revisions, and has practical implications for nosological policy and clinical education."
)

KEYWORDS = [
    "Sapir–Whorf hypothesis",
    "nosology",
    "disease classification",
    "linguistic relativity",
    "looping effects",
    "philosophy of medicine",
    "ICD",
    "diagnostic labels",
]

# References (Vancouver style, numbered in order of appearance)
REFERENCES = [
    # 1
    "Sapir E. The status of linguistics as a science. Language. 1929;5(4):207–14.",
    # 2
    "Whorf BL. Language, thought, and reality: selected writings of Benjamin Lee Whorf. Carroll JB, editor. Cambridge (MA): MIT Press; 1956.",
    # 3
    "Warner R. The relationship between language and disease concepts. Int J Psychiatry Med. 1976;7(1):57–68.",
    # 4
    "Hacking I. The looping effects of human kinds. In: Sperber D, Premack D, Premack AJ, editors. Causal cognition: a multidisciplinary debate. Oxford: Clarendon Press; 1995. p. 351–94.",
    # 5
    "Hacking I. Making up people. London Review of Books. 2006;28(16):23–6.",
    # 6
    "Kleinman A. The illness narratives: suffering, healing, and the human condition. New York: Basic Books; 1988.",
    # 7
    "Jutel A. Sociology of diagnosis: a preliminary review. Sociol Health Illn. 2009;31(2):278–99.",
    # 8
    "Zachar P, Kendler KS. The philosophy of nosology. Annu Rev Clin Psychol. 2017;13:49–71.",
    # 9
    "Boorse C. Health as a theoretical concept. Philos Sci. 1977;44(4):542–73.",
    # 10
    "Swartz L. Anorexia nervosa as a culture-bound syndrome. Soc Sci Med. 1985;20(7):725–30.",
    # 11
    "Michaleff ZA, Glasziou P, Thomas R. Consequences of a diagnostic label: a systematic scoping review and thematic framework. Front Public Health. 2021;9:725877.",
    # 12
    "Nickel B, Moynihan R, Barratt A, Brito JP, McCaffery K. Words do matter: a systematic review on how different terminology for the same condition influences management preferences. BMJ Open. 2017;7(7):e014129.",
    # 13
    "Iwasaki K, Takahashi M, Nakata A. Health problems due to long working hours in Japan: working hours, workers' compensation (karoshi), and preventive measures. Ind Health. 2006;44(4):537–40.",
    # 14
    "Nishiyama K, Johnson JV. Karoshi—death from overwork: occupational health consequences of Japanese production management. Int J Health Serv. 1997;27(4):625–41.",
    # 15
    "World Health Organization. ICD-11 for Mortality and Morbidity Statistics. Geneva: WHO; 2019.",
    # 16
    "Reed GM, First MB, Kogan CS, et al. Innovations and changes in the ICD-11 classification of mental, behavioural and neurodevelopmental disorders. World Psychiatry. 2019;18(1):3–19.",
    # 17
    "Treede RD, Rief W, Barke A, et al. Chronic pain as a symptom or a disease: the IASP Classification of Chronic Pain for the International Classification of Diseases (ICD-11). Pain. 2019;160(1):19–27.",
    # 18
    "Tsou JY. Natural kinds, psychiatric classification and the history of the DSM. Hist Psychiatry. 2016;27(4):406–24.",
    # 19
    "Cooper R. Classifying madness: a philosophical examination of the Diagnostic and Statistical Manual of Mental Disorders. Dordrecht: Springer; 2005.",
    # 20
    "Fabrega H Jr. Disease and social behavior: an interdisciplinary perspective. Cambridge (MA): MIT Press; 1974.",
    # 21
    "Eisenberg L. Disease and illness: distinctions between professional and popular ideas of sickness. Cult Med Psychiatry. 1977;1(1):9–23.",
    # 22
    "Rosenhan DL. On being sane in insane places. Science. 1973;179(4070):250–8.",
    # 23
    "Thibault JM, Bhatt DL, Engel GL. The biopsychosocial model: past, present, future. Psychosomatics. 2003;44(4):267–75.",
    # 24
    "Boroditsky L. Does language shape thought? Mandarin and English speakers' conceptions of time. Cogn Psychol. 2001;43(1):1–22.",
    # 25
    "Lupyan G, Bergen B. How language programs the mind. Top Cogn Sci. 2016;8(2):408–24.",
    # 26
    "Conrad P. The medicalization of society: on the transformation of human conditions into treatable disorders. Baltimore: Johns Hopkins University Press; 2007.",
    # 27
    "Bowker GC, Star SL. Sorting things out: classification and its consequences. Cambridge (MA): MIT Press; 1999.",
    # 28
    "Craddock N, Owen MJ. The Kraepelinian dichotomy — going, going... but still not gone. Br J Psychiatry. 2010;196(2):92–5.",
    # 29
    "Frances A. Saving normal: an insider's revolt against out-of-control psychiatric diagnosis, DSM-5, Big Pharma, and the medicalization of ordinary life. New York: William Morrow; 2013.",
    # 30
    "Kirmayer LJ. Cultural variations in the clinical presentation of depression and anxiety: implications for diagnosis and treatment. J Clin Psychiatry. 2001;62 Suppl 13:22–8.",
]

# ---------------------------------------------------------------------------
# BODY TEXT sections  (with {n} citation markers for superscript)
# ---------------------------------------------------------------------------

INTRO_PARAS = [
    (
        "The Sapir–Whorf hypothesis, formulated in the early twentieth century by Edward Sapir and "
        "Benjamin Lee Whorf, proposes that the structure of a language influences its speakers' "
        "cognition and worldview.{1,2} In its strong form (linguistic determinism), language "
        "determines thought; in its weak form (linguistic relativity), language influences habitual "
        "patterns of thought without fully constraining them.{24,25} Although the strong form has "
        "been largely abandoned in linguistics, robust evidence supports the weak form across "
        "domains including color perception, spatial reasoning, and temporal cognition.{24}"
    ),
    (
        "An analogous phenomenon may operate in medicine. Nosological frameworks—the formal systems "
        "by which diseases are classified, named, and organized—constitute the 'language' of "
        "clinical medicine. Just as natural language categories shape how speakers perceive and "
        "reason about the world, nosological categories may shape how clinicians perceive, diagnose, "
        "and treat patients. Warner first explored this analogy in 1976, arguing that linguistic "
        "structures in different cultures lead to fundamentally different conceptions of disease.{3} "
        "However, the parallel has not been formally developed beyond initial observations."
    ),
    (
        "The claim that 'disease concepts shape medical practice' might appear trivially true—medicine "
        "is, after all, organized around disease categories. However, we argue that the effects "
        "extend well beyond the obvious and warrant systematic theoretical treatment. The phenomenon "
        "includes cognitive effects on clinical reasoning (diagnostic foreclosure), institutional "
        "effects on research funding and treatment guidelines, and population-level effects on "
        "epidemiological patterns. Moreover, as Ian Hacking has demonstrated for psychiatric "
        "classifications, a distinctive looping mechanism operates: classification systems alter "
        "patient self-identification and symptom expression, which in turn generates data that "
        "reinforces the original classification.{4,5}"
    ),
    (
        "At the same time, evidence constrains this hypothesis. Kleinman's observation that "
        "anorexia nervosa occurs even in cultures where thinness is not idealized{6,10} suggests "
        "that some diseases possess a biological substrate that manifests independently of "
        "nosological framing. Any adequate theoretical framework must accommodate both the "
        "constructive power of nosological categories and the existence of nosology-independent "
        "disease phenomena."
    ),
    (
        "In this paper, we propose the 'Nosological Relativity' (NR) framework—a formal analogue "
        "of the Sapir–Whorf hypothesis applied to medical classification systems. We define the "
        "framework's core propositions, distinguish strong and weak forms, delineate three levels "
        "of effect, and formalize the looping mechanism as the central feedback dynamic. We then "
        "illustrate the framework with case studies and derive testable predictions suitable for "
        "empirical investigation."
    ),
]

BACKGROUND_PARAS = {
    "The Sapir–Whorf Hypothesis in Linguistics": [
        (
            "The linguistic relativity hypothesis has undergone substantial revision since its "
            "original formulation.{1,2} Contemporary research supports a moderate position: "
            "language does not rigidly determine thought, but it does influence habitual patterns "
            "of cognition.{24,25} Boroditsky demonstrated that Mandarin and English speakers "
            "conceptualize time differently, consistent with structural differences in how each "
            "language encodes temporal relations.{24} Lupyan and Bergen showed that linguistic "
            "labels facilitate categorization and modulate perceptual processing.{25} These "
            "findings establish that symbolic classification systems exert measurable cognitive "
            "effects—a principle we extend to medical nosology."
        ),
    ],
    "Prior Work on Language and Disease": [
        (
            "Warner's 1976 paper represents the first explicit application of the Sapir–Whorf "
            "hypothesis to medicine.{3} He argued that Indo-European linguistic structures—"
            "particularly the use of nouns rather than verbs to describe illness, the extensive "
            "use of spatial metaphors, and the subject–predicate dichotomy—encourage a static, "
            "unicausal, and dualistic conception of disease. Warner suggested that these linguistic "
            "features promote over-reliance on surgery and inhibit recognition of social and "
            "psychological factors in disease causation."
        ),
        (
            "Subsequent work has expanded this perspective without formalization. Eisenberg "
            "distinguished between 'disease' (the biomedical construct) and 'illness' (the "
            "patient's lived experience), showing that the gap between these constructs is "
            "mediated by cultural and linguistic categories.{21} Fabrega offered an "
            "interdisciplinary framework linking disease to social behavior.{20} Kleinman's "
            "medical anthropology demonstrated that cultural categories of sickness guide "
            "labeling, help-seeking, and therapeutic responses.{6} Hacking's philosophy of "
            "human kinds introduced the looping-effect concept, showing that psychiatric "
            "classifications are not passive labels but active constituents of the phenomena "
            "they describe.{4,5} Jutel's sociology of diagnosis showed that the diagnostic "
            "act itself is a social event with consequences that extend beyond the clinical "
            "encounter.{7}"
        ),
        (
            "Despite these rich intellectual traditions, no integrated formal framework exists "
            "for the proposition that nosological categories shape medical practice in ways "
            "analogous to linguistic relativity. The present paper fills this gap."
        ),
    ],
}

FRAMEWORK_INTRO = (
    "We now present the Nosological Relativity (NR) framework. The framework consists of "
    "a core thesis, six formal propositions, three levels of effect, and a looping-mechanism "
    "model (Fig. 1)."
)

PROPOSITIONS = [
    (
        "Proposition 1 (Core Thesis).",
        "Nosological frameworks—the systems by which diseases are classified and named—"
        "do not merely describe pre-existing clinical reality but actively shape the "
        "perception, categorization, and management of illness at multiple levels of "
        "the healthcare system."
    ),
    (
        "Proposition 2 (Weak Form).",
        "Nosological categories influence but do not fully determine clinical reasoning, "
        "research priorities, and patient outcomes. Clinicians can reason outside the "
        "categories of their training, but such reasoning requires additional cognitive "
        "effort and is statistically less frequent."
    ),
    (
        "Proposition 3 (Strong Form).",
        "For certain conditions—particularly those lacking objective biomarkers—nosological "
        "categories effectively determine clinical reality: a condition that does not exist "
        "in the classification system does not receive diagnosis, treatment, or research "
        "attention, and is therefore functionally non-existent within the healthcare system."
    ),
    (
        "Proposition 4 (Cognitive Level).",
        "Nosological categories function as cognitive schemata for clinicians, producing "
        "diagnostic foreclosure (cessation of diagnostic reasoning once a category is "
        "assigned), anchoring bias toward category-consistent features, and "
        "difficulty perceiving cross-category patterns."
    ),
    (
        "Proposition 5 (Institutional Level).",
        "Nosological categories structure research funding allocation, clinical guideline "
        "development, specialist training, and reimbursement systems. Conditions without "
        "nosological representation receive disproportionately less institutional support "
        "relative to their disease burden."
    ),
    (
        "Proposition 6 (Population Level).",
        "Nosological categories shape epidemiological patterns through differential "
        "case-finding, diagnostic migration (reclassification of existing cases when "
        "categories change), and the creation of new illness identities that alter "
        "help-seeking behavior."
    ),
]

LOOPING_PARAS = [
    (
        "Central to the NR framework is the looping mechanism, adapted from Hacking's "
        "analysis of human kinds (Fig. 2).{4,5} The mechanism operates as follows:"
    ),
    (
        "Stage 1 (Classification → Practice): A nosological category is introduced or "
        "revised within a formal system (e.g., ICD, DSM). This creates diagnostic criteria, "
        "treatment algorithms, and referral pathways organized around the category."
    ),
    (
        "Stage 2 (Practice → Experience): Clinicians using the category communicate the "
        "diagnosis to patients, who incorporate it into their self-understanding. The "
        "diagnostic label provides a framework for interpreting symptoms, a basis for "
        "illness identity, and a social role.{11}"
    ),
    (
        "Stage 3 (Experience → Data): Patients who self-identify with the diagnostic "
        "category generate health-seeking behavior, symptom reports, and clinical "
        "encounters that are coded using the category. This produces epidemiological "
        "data consistent with the category's existence.{26}"
    ),
    (
        "Stage 4 (Data → Classification): Epidemiological data are used to validate and "
        "refine nosological categories. Categories that generate data consistent with "
        "their definitions appear empirically supported, completing the loop."
    ),
    (
        "This looping mechanism creates a self-reinforcing cycle that can make it "
        "difficult to distinguish between categories that 'carve nature at its joints' "
        "and categories that generate the very patterns they claim to describe.{18,19}"
    ),
]

EVIDENCE_SECTIONS = {
    "Diagnostic Foreclosure in Pain Medicine": [
        (
            "The recognition of chronic pain as an independent disease entity in ICD-11 "
            "(code MG30) illustrates the cognitive-level effects of nosological categories.{17} "
            "Prior to ICD-11, chronic pain was classified exclusively as a symptom of underlying "
            "conditions. This nosological framing constrained clinical reasoning: once an "
            "underlying condition was identified (e.g., osteoarthritis), diagnostic exploration "
            "of pain mechanisms typically ceased—a phenomenon we term diagnostic foreclosure.{4}"
        ),
        (
            "The introduction of chronic pain as an independent diagnostic entity creates the "
            "conditions for a natural experiment. The NR framework predicts that following "
            "adoption of ICD-11's chronic pain categories, clinicians will: (a) refer more "
            "patients to pain specialists rather than treating pain exclusively through the "
            "lens of underlying conditions; (b) increase prescribing of pain-specific therapies; "
            "and (c) generate new epidemiological data on chronic pain prevalence that will be "
            "attributed to improved detection rather than to the nosological change itself."
        ),
        (
            "Nickel et al. showed in a systematic review that different terminology for the "
            "same condition influences clinicians' and patients' management preferences.{12} "
            "This finding directly supports Proposition 2: nosological labels influence, but do "
            "not determine, clinical reasoning."
        ),
    ],
    "Karoshi: Culture-Specific Nosological Innovation": [
        (
            "Karoshi (過労死, death from overwork) provides a compelling illustration of "
            "Proposition 3—the strong form of nosological relativity. The concept emerged in "
            "Japan in the 1970s as a legally and medically recognized cause of death, distinct "
            "from the underlying cardiovascular or cerebrovascular events.{13,14} In other "
            "countries, identical pathophysiological events (myocardial infarction, stroke) "
            "occurring in the context of overwork are classified under their organ-specific "
            "categories without reference to occupational causation."
        ),
        (
            "The existence of karoshi as a nosological category in Japan has produced a "
            "cascade of effects consistent with the NR framework: (a) at the cognitive level, "
            "Japanese clinicians consider occupational overwork in the differential diagnosis "
            "of sudden cardiovascular death in working-age adults, a consideration that is "
            "largely absent in countries without an equivalent category; (b) at the "
            "institutional level, Japan has developed a dedicated legal-medical apparatus for "
            "karoshi certification, with specific overtime-hour thresholds (the '80-hour "
            "karoshi line') that do not exist elsewhere; (c) at the population level, Japan "
            "reports karoshi statistics as a distinct epidemiological phenomenon, generating "
            "policy responses (the 2014 Act on Promotion of Measures to Prevent Karoshi) "
            "that are structurally impossible in countries where the category does not exist."
        ),
        (
            "The ICD-11 inclusion of burnout (QD85) as a classifiable occupational phenomenon "
            "provides a partial international diffusion of this nosological innovation, and "
            "the NR framework predicts that countries adopting this code will begin to "
            "replicate—at a smaller scale—the institutional and epidemiological patterns "
            "observed in Japan.{15,16}"
        ),
    ],
    "ICD Revisions as Natural Experiments": [
        (
            "Major revisions of the International Classification of Diseases constitute "
            "natural experiments for the NR framework. The transition from ICD-10 to ICD-11 "
            "introduced several nosological innovations whose effects can be tracked using "
            "before-after designs (Table 1).{15,16}"
        ),
        (
            "Of particular theoretical interest is the reclassification of gender incongruence "
            "from Mental and Behavioural Disorders (ICD-10, Chapter V) to Conditions Related "
            "to Sexual Health (ICD-11, Chapter 17). The NR framework predicts that this "
            "reclassification will produce measurable changes at all three levels: reduced "
            "cognitive association with mental illness among clinicians, institutional "
            "reorganization of care pathways away from psychiatric services, and population-level "
            "reduction in stigma-related barriers to care.{16}"
        ),
    ],
    "Counter-Evidence: Disease-Independent-of-Nosology": [
        (
            "The NR framework must accommodate evidence that some diseases manifest independently "
            "of nosological framing. Kleinman documented that anorexia nervosa occurs in "
            "cultures where thinness is not idealized,{6} and Swartz argued that anorexia should "
            "be understood as a culture-bound syndrome precisely because its form varies across "
            "cultures even when its core features persist.{10} More broadly, infectious diseases "
            "with clear etiological agents (e.g., tuberculosis, malaria) manifest regardless "
            "of how they are classified."
        ),
        (
            "These observations delineate the boundary conditions of nosological relativity. "
            "Boorse's biostatistical theory of health provides a useful reference point: "
            "to the extent that a condition deviates from species-typical function in ways "
            "measurable independently of nosological framing, it is less susceptible to "
            "nosological relativity.{9} Kirmayer has further documented how cultural "
            "variations in the clinical presentation of depression and anxiety complicate "
            "simple universalist assumptions.{30} We propose that the strength of "
            "nosological effects is inversely proportional to the strength of the underlying "
            "biological signal: conditions with strong, objectively measurable biological "
            "substrates (e.g., fractures, infections with identifiable pathogens) are less "
            "susceptible to nosological shaping, while conditions defined primarily by "
            "symptom clusters, functional impairment, or subjective experience (e.g., "
            "fibromyalgia, chronic fatigue syndrome, many psychiatric disorders) are highly "
            "susceptible. This gradient constitutes a testable prediction of the NR framework."
        ),
    ],
}

PREDICTIONS_PARAS = [
    (
        "The NR framework generates three empirically testable predictions:"
    ),
    (
        "Prediction 1 (ICD Revision Effect): Introduction of a new diagnostic category "
        "in ICD-11 will be followed by increased detection rates, specialist referrals, "
        "and targeted treatment prescriptions for the corresponding condition, independent "
        "of any change in the underlying biological incidence. This can be tested using "
        "interrupted time-series analysis of claims data around ICD transition dates."
    ),
    (
        "Prediction 2 (Biomarker Gradient): The magnitude of nosological effects on "
        "diagnostic and treatment patterns will be inversely correlated with the "
        "availability and specificity of objective biomarkers for the condition. Conditions "
        "diagnosed purely on clinical criteria will show larger effects than those with "
        "definitive laboratory or imaging markers."
    ),
    (
        "Prediction 3 (Cross-National Divergence): Countries with different nosological "
        "traditions for the same pathophysiological phenomenon (e.g., Japan's karoshi "
        "versus other countries' organ-specific coding of overwork-related cardiovascular "
        "death) will show systematically different patterns of diagnosis, treatment, "
        "occupational health policy, and patient self-identification for those conditions."
    ),
]

DISCUSSION_PARAS = [
    (
        "The Nosological Relativity framework offers several implications for medical "
        "practice and nosological policy."
    ),
    (
        "First, nosological revisions should be understood not merely as improved "
        "descriptions of clinical reality but as interventions that actively reshape "
        "clinical practice. This suggests that the process of revising classification "
        "systems should incorporate prospective impact assessment, analogous to "
        "regulatory impact assessment for legislation.{8,27}"
    ),
    (
        "Second, clinical education should include explicit instruction in the cognitive "
        "effects of diagnostic categories. Awareness of diagnostic foreclosure—the "
        "tendency to cease diagnostic reasoning once a category is assigned—may partially "
        "mitigate its effects, analogous to debiasing training for other cognitive "
        "heuristics.{22} Integration with biopsychosocial approaches may further "
        "counteract the reductionism inherent in categorical thinking.{23}"
    ),
    (
        "Third, the framework provides a principled basis for evaluating proposals to "
        "create new diagnostic categories. The benefits of nosological recognition "
        "(improved detection, research funding, treatment development) must be weighed "
        "against the risks of reification (treating a provisional category as a natural "
        "kind) and looping effects (the category generating the very phenomenon it "
        "describes).{26,29}"
    ),
    (
        "Limitations of this framework should be acknowledged. The NR framework is "
        "primarily a theoretical construct at this stage; the testable predictions "
        "outlined above require empirical validation. The analogy between natural "
        "language and nosological systems is imperfect: natural languages are acquired "
        "implicitly and used by entire populations, while nosological systems are "
        "explicitly constructed and primarily used by specialists. Additionally, the "
        "framework does not address the normative question of whether nosological "
        "categories should be designed to minimize or to harness their constructive "
        "effects.{8}"
    ),
]

CONCLUSION_PARAS = [
    (
        "We have proposed the Nosological Relativity framework as a formal theoretical "
        "lens for understanding how disease classification systems shape medical practice. "
        "Drawing on the Sapir–Whorf hypothesis from linguistics, Hacking's looping-effect "
        "mechanism from philosophy of science, and evidence from medical anthropology and "
        "the sociology of diagnosis, we have defined six propositions, three levels of "
        "effect, and three testable predictions."
    ),
    (
        "The framework does not deny the existence of biological disease substrates; "
        "rather, it argues that the nosological 'lens' through which these substrates "
        "are perceived, categorized, and managed introduces systematic and predictable "
        "distortions. By making these distortions explicit and testable, the NR framework "
        "opens new avenues for empirical research and provides principled guidance for "
        "nosological policy.{27,28}"
    ),
    (
        "Future research should prioritize empirical testing of the three predictions "
        "outlined above, particularly using the ICD-10 to ICD-11 transition as a "
        "natural experiment. Cross-national comparative studies of conditions with "
        "divergent nosological traditions (such as karoshi) offer additional promising "
        "research designs. Ultimately, the Nosological Relativity framework invites a "
        "reflexive turn in medicine: the recognition that our systems for classifying "
        "disease are not neutral mirrors of biological reality, but active participants "
        "in the construction of clinical knowledge."
    ),
]

# Table 1 data
TABLE1_CAPTION = "Table 1. Selected ICD-11 Nosological Innovations and Predicted Effects Under the Nosological Relativity Framework"
TABLE1_HEADERS = ["ICD-11 Change", "Cognitive Level", "Institutional Level", "Population Level"]
TABLE1_ROWS = [
    [
        "Chronic pain as independent entity (MG30)",
        "Clinicians consider pain as primary diagnosis, not merely symptom",
        "Pain medicine receives dedicated funding streams and guidelines",
        "Increased chronic pain prevalence rates attributed to improved detection",
    ],
    [
        "Burnout as occupational phenomenon (QD85)",
        "Clinicians screen for occupational factors in exhaustion presentations",
        "Occupational health services develop burnout-specific protocols",
        "New epidemiological data on burnout prevalence emerge",
    ],
    [
        "Gaming disorder (6C51)",
        "Clinicians identify pathological gaming as a diagnosable condition",
        "Treatment centers and clinical guidelines for gaming disorder develop",
        "Incidence data for gaming disorder are generated and reported",
    ],
    [
        "Complex PTSD (6B41)",
        "Clinicians distinguish complex from simple PTSD in trauma patients",
        "Research funding specifically allocated to complex PTSD mechanisms",
        "Existing PTSD patients reclassified, altering prevalence estimates",
    ],
    [
        "Gender incongruence reclassification (HA60–HA6Z)",
        "Reduced cognitive association with mental illness",
        "Care pathways reorganized away from psychiatric services",
        "Reduced stigma-related barriers to care access",
    ],
]


# ---------------------------------------------------------------------------
# Build the manuscript DOCX
# ---------------------------------------------------------------------------

def build_manuscript(fig1_path, fig2_path):
    doc = Document()

    # --- Style setup ---
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Times New Roman'
    font.size = Pt(11)
    style.paragraph_format.line_spacing = 1.5

    # ---- TITLE PAGE ----
    tp = doc.add_paragraph()
    tp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = tp.add_run(TITLE)
    run.bold = True
    run.font.size = Pt(16)
    run.font.name = 'Times New Roman'

    tp2 = doc.add_paragraph()
    tp2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run2 = tp2.add_run(AUTHORS)
    run2.font.size = Pt(12)

    if AFFILIATIONS:
        tp3 = doc.add_paragraph()
        tp3.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run3 = tp3.add_run(AFFILIATIONS)
        run3.font.size = Pt(10)
        run3.italic = True

    doc.add_paragraph()  # spacer

    # ---- ABSTRACT ----
    add_heading(doc, 'Abstract', level=1)
    add_body(doc, ABSTRACT)

    kw_para = doc.add_paragraph()
    kw_run = kw_para.add_run('Keywords: ')
    kw_run.bold = True
    kw_run.font.size = Pt(11)
    kw_text = kw_para.add_run('; '.join(KEYWORDS))
    kw_text.font.size = Pt(11)

    doc.add_page_break()

    # ---- INTRODUCTION ----
    add_heading(doc, '1. Introduction', level=1)
    for para_text in INTRO_PARAS:
        add_body(doc, para_text)

    # ---- BACKGROUND ----
    add_heading(doc, '2. Background', level=1)
    section_num = 1
    for sub_title, paras in BACKGROUND_PARAS.items():
        add_heading(doc, f'2.{section_num} {sub_title}', level=2)
        section_num += 1
        for para_text in paras:
            add_body(doc, para_text)



    # ---- THE NOSOLOGICAL RELATIVITY FRAMEWORK ----
    add_heading(doc, '3. The Nosological Relativity Framework', level=1)
    add_body(doc, FRAMEWORK_INTRO)

    # --- Figure 1 inline ---
    doc.add_paragraph()  # spacer
    fig1_para = doc.add_paragraph()
    fig1_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = fig1_para.add_run()
    run.add_picture(str(fig1_path), width=Inches(5.5))

    cap1 = doc.add_paragraph()
    cap1.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cap1.paragraph_format.space_before = Pt(12)
    cap_run = cap1.add_run(
        'Fig. 1. The Nosological Relativity Framework. '
        'Nosological categories shape medical practice at three levels '
        '(cognitive, institutional, population) via the looping mechanism. '
        'The strength of effect varies along a spectrum from weak (influence) '
        'to strong (determination). Counter-evidence from nosology-independent '
        'diseases delineates the framework\'s boundary conditions.'
    )
    cap_run.italic = True
    cap_run.font.size = Pt(10)
    doc.add_paragraph()  # spacer

    # --- Propositions ---
    add_heading(doc, '3.1 Core Propositions', level=2)
    for prop_title, prop_text in PROPOSITIONS:
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(6)
        p.paragraph_format.line_spacing = 1.5
        run_t = p.add_run(prop_title + ' ')
        run_t.bold = True
        run_t.font.size = Pt(11)
        run_b = p.add_run(prop_text)
        run_b.font.size = Pt(11)

    # --- Looping mechanism ---
    add_heading(doc, '3.2 The Looping Mechanism', level=2)
    for para_text in LOOPING_PARAS:
        add_body(doc, para_text)

    # --- Figure 2 inline ---
    doc.add_paragraph()  # spacer
    fig2_para = doc.add_paragraph()
    fig2_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = fig2_para.add_run()
    run.add_picture(str(fig2_path), width=Inches(5.5))

    cap2 = doc.add_paragraph()
    cap2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cap2.paragraph_format.space_before = Pt(12)
    cap_run2 = cap2.add_run(
        'Fig. 2. The Nosological Looping Mechanism. '
        'A four-stage cycle in which nosological classification shapes clinical practice, '
        'which molds patient experience, which generates epidemiological data, '
        'which in turn revises the classification. This self-reinforcing loop is the '
        'central dynamic of the Nosological Relativity framework.'
    )
    cap_run2.italic = True
    cap_run2.font.size = Pt(10)
    doc.add_paragraph()  # spacer

    # ---- EVIDENCE AND CASE STUDIES ----
    add_heading(doc, '4. Evidence and Case Studies', level=1)
    section_num = 1
    for sub_title, paras in EVIDENCE_SECTIONS.items():
        add_heading(doc, f'4.{section_num} {sub_title}', level=2)
        section_num += 1
        for i, para_text in enumerate(paras):
            add_body(doc, para_text)
            # Insert Table 1 after the first paragraph of ICD Revisions section
            if 'ICD Revisions' in sub_title and i == 0:
                # Table 1
                doc.add_paragraph()  # spacer before table
                t1cap = doc.add_paragraph()
                t1cap.paragraph_format.space_before = Pt(14)
                t1_run = t1cap.add_run(TABLE1_CAPTION)
                t1_run.bold = True
                t1_run.font.size = Pt(10)

                table = doc.add_table(rows=1 + len(TABLE1_ROWS), cols=4)
                table.style = 'Table Grid'
                table.alignment = WD_TABLE_ALIGNMENT.CENTER

                # Headers
                for j, header in enumerate(TABLE1_HEADERS):
                    set_cell_text(table.rows[0].cells[j], header, bold=True, size=Pt(9))

                # Rows
                for r_idx, row_data in enumerate(TABLE1_ROWS):
                    for c_idx, cell_text in enumerate(row_data):
                        set_cell_text(table.rows[r_idx + 1].cells[c_idx], cell_text, size=Pt(9))

                # Set column widths
                widths = [Inches(1.6), Inches(1.6), Inches(1.6), Inches(1.6)]
                for row in table.rows:
                    for idx, width in enumerate(widths):
                        row.cells[idx].width = width

                doc.add_paragraph()  # spacer after table

    # ---- TESTABLE PREDICTIONS ----
    add_heading(doc, '5. Testable Predictions', level=1)
    for para_text in PREDICTIONS_PARAS:
        add_body(doc, para_text)

    # ---- DISCUSSION ----
    add_heading(doc, '6. Discussion', level=1)
    for para_text in DISCUSSION_PARAS:
        add_body(doc, para_text)

    # ---- CONCLUSION ----
    add_heading(doc, '7. Conclusion', level=1)
    for para_text in CONCLUSION_PARAS:
        add_body(doc, para_text)

    # ---- REFERENCES ----
    doc.add_page_break()
    add_heading(doc, 'References', level=1)
    for i, ref in enumerate(REFERENCES, 1):
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(2)
        p.paragraph_format.line_spacing = 1.15
        run_num = p.add_run(f'{i}. ')
        run_num.bold = True
        run_num.font.size = Pt(10)
        run_ref = p.add_run(ref)
        run_ref.font.size = Pt(10)

    # Save
    out_path = OUT_DIR / "manuscript_en.docx"
    doc.save(str(out_path))
    print(f"Manuscript saved: {out_path}")
    return out_path


# ---------------------------------------------------------------------------
# Build editable figures PPTX
# ---------------------------------------------------------------------------

def build_figures_pptx(fig1_path, fig2_path):
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    for fig_path, title_text, caption_text in [
        (fig1_path,
         "Fig. 1. The Nosological Relativity Framework",
         "Nosological categories shape medical practice at three levels "
         "(cognitive, institutional, population) via the looping mechanism."),
        (fig2_path,
         "Fig. 2. The Nosological Looping Mechanism",
         "A four-stage cycle: classification → practice → experience → data → classification."),
    ]:
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)

        # Title
        txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3),
                                          PptxInches(12.333), PptxInches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = PptxPt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Image
        img_left = PptxInches(1.5)
        img_top = PptxInches(1.3)
        img_width = PptxInches(10.333)
        slide.shapes.add_picture(str(fig_path), img_left, img_top, width=img_width)

        # Caption
        capBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(6.3),
                                           PptxInches(12.333), PptxInches(1.0))
        cf = capBox.text_frame
        cf.word_wrap = True
        cp = cf.paragraphs[0]
        cp.text = caption_text
        cp.font.size = PptxPt(14)
        cp.font.italic = True
        cp.alignment = PP_ALIGN.CENTER

    out_path = OUT_DIR / "figures_en.pptx"
    prs.save(str(out_path))
    print(f"Figures PPTX saved: {out_path}")
    return out_path


# ---------------------------------------------------------------------------
# Build editable tables DOCX
# ---------------------------------------------------------------------------

def build_tables_docx():
    doc = Document()
    style = doc.styles['Normal']
    style.font.name = 'Times New Roman'
    style.font.size = Pt(10)

    add_heading(doc, TABLE1_CAPTION, level=2)

    table = doc.add_table(rows=1 + len(TABLE1_ROWS), cols=4)
    table.style = 'Table Grid'

    for j, header in enumerate(TABLE1_HEADERS):
        set_cell_text(table.rows[0].cells[j], header, bold=True, size=Pt(10))

    for r_idx, row_data in enumerate(TABLE1_ROWS):
        for c_idx, cell_text in enumerate(row_data):
            set_cell_text(table.rows[r_idx + 1].cells[c_idx], cell_text, size=Pt(10))

    out_path = OUT_DIR / "tables_en.docx"
    doc.save(str(out_path))
    print(f"Tables DOCX saved: {out_path}")
    return out_path


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print("Generating figures...")
    fig1 = create_figure1()
    fig2 = create_figure2()

    print("Building manuscript...")
    ms = build_manuscript(fig1, fig2)

    print("Building figures PPTX...")
    build_figures_pptx(fig1, fig2)

    print("Building tables DOCX...")
    build_tables_docx()

    print("Done!")


if __name__ == '__main__':
    main()
