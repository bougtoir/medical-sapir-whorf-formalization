#!/usr/bin/env python3
"""
Generate editable PPTX for figures and tables.
One slide per figure/table, with title and caption.

Outputs:
  medical_sapir_whorf/output/figures_tables_en.pptx
  medical_sapir_whorf/output/figures_tables_ja.pptx
"""

import os
import io
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.font_manager as fm

# Ensure CJK font is available for Japanese figures
_CJK_FONT = None
for _f in fm.fontManager.ttflist:
    if "Noto Sans CJK JP" in _f.name and "Bold" not in _f.fname:
        _CJK_FONT = _f.fname
        break
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

OUT_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "output")
os.makedirs(OUT_DIR, exist_ok=True)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def create_figure1_en():
    """Disease-Illness-Sickness feedback model (English)."""
    fig, ax = plt.subplots(1, 1, figsize=(10, 6))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 7)
    ax.axis("off")

    # Three boxes
    boxes = [
        (1.5, 5.0, "DISEASE\n(Biomedical pathology)", "#4ECDC4"),
        (4.0, 2.0, "SICKNESS\n(Social/institutional\nrecognition)", "#FF6B6B"),
        (6.5, 5.0, "ILLNESS\n(Subjective patient\nexperience)", "#45B7D1"),
    ]
    for x, y, txt, color in boxes:
        rect = mpatches.FancyBboxPatch(
            (x, y), 2.5, 1.5, boxstyle="round,pad=0.15",
            facecolor=color, edgecolor="black", linewidth=2, alpha=0.85
        )
        ax.add_patch(rect)
        ax.text(x + 1.25, y + 0.75, txt, ha="center", va="center",
                fontsize=10, fontweight="bold", color="white")

    # Nosological categories label
    rect2 = mpatches.FancyBboxPatch(
        (3.5, 0.3), 3.5, 0.8, boxstyle="round,pad=0.1",
        facecolor="#FFD93D", edgecolor="black", linewidth=2
    )
    ax.add_patch(rect2)
    ax.text(5.25, 0.7, "NOSOLOGICAL\nCATEGORIES (ICD, DSM)",
            ha="center", va="center", fontsize=9, fontweight="bold")

    # Arrows
    arrow_props = dict(arrowstyle="->,head_width=0.3,head_length=0.2",
                       color="black", linewidth=2)

    # Nosological -> Sickness
    ax.annotate("", xy=(5.25, 2.0), xytext=(5.25, 1.1),
                arrowprops=arrow_props)
    ax.text(5.7, 1.5, "Defines\ncategories", fontsize=8, color="gray")

    # Sickness -> Disease (feedback)
    ax.annotate("", xy=(2.75, 5.0), xytext=(4.0, 3.5),
                arrowprops=dict(arrowstyle="->", color="#FF6B6B",
                                linewidth=2, linestyle="--"))
    ax.text(2.2, 4.2, "Shapes diagnosis\n& treatment", fontsize=8,
            color="#FF6B6B", ha="center")

    # Sickness -> Illness (feedback)
    ax.annotate("", xy=(7.75, 5.0), xytext=(6.5, 3.5),
                arrowprops=dict(arrowstyle="->", color="#FF6B6B",
                                linewidth=2, linestyle="--"))
    ax.text(8.2, 4.2, "Shapes patient\nself-understanding",
            fontsize=8, color="#FF6B6B", ha="center")

    # Disease <-> Illness (traditional)
    ax.annotate("", xy=(6.5, 5.75), xytext=(4.0, 5.75),
                arrowprops=dict(arrowstyle="<->", color="gray",
                                linewidth=1.5))
    ax.text(5.25, 6.1, "Traditional relationship", fontsize=8,
            color="gray", ha="center")

    # Sapir-Whorf label
    ax.text(5.25, 6.7, "Medical Sapir\u2013Whorf Feedback Loop",
            ha="center", va="center", fontsize=14, fontweight="bold",
            color="#2C3E50")

    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight",
                facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def create_figure1_ja():
    """Disease-Illness-Sickness feedback model (Japanese)."""
    plt.rcParams["font.family"] = "Noto Sans CJK JP"
    fig, ax = plt.subplots(1, 1, figsize=(10, 6))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 7)
    ax.axis("off")

    boxes = [
        (1.5, 5.0, "DISEASE\n(\u75be\u60a3\uff1a\u751f\u7269\u533b\u5b66\u7684\u75c5\u614b)", "#4ECDC4"),
        (4.0, 2.0, "SICKNESS\n(\u75c5\u4eba\u5f79\u5272\uff1a\u793e\u4f1a\u7684\u30fb\n\u5236\u5ea6\u7684\u8a8d\u77e5)", "#FF6B6B"),
        (6.5, 5.0, "ILLNESS\n(\u75c5\u3044\uff1a\u60a3\u8005\u306e\n\u4e3b\u89b3\u7684\u4f53\u9a13)", "#45B7D1"),
    ]
    for x, y, txt, color in boxes:
        rect = mpatches.FancyBboxPatch(
            (x, y), 2.5, 1.5, boxstyle="round,pad=0.15",
            facecolor=color, edgecolor="black", linewidth=2, alpha=0.85
        )
        ax.add_patch(rect)
        ax.text(x + 1.25, y + 0.75, txt, ha="center", va="center",
                fontsize=10, fontweight="bold", color="white")

    rect2 = mpatches.FancyBboxPatch(
        (3.2, 0.3), 4.1, 0.8, boxstyle="round,pad=0.1",
        facecolor="#FFD93D", edgecolor="black", linewidth=2
    )
    ax.add_patch(rect2)
    ax.text(5.25, 0.7, "\u75be\u75c5\u5206\u985e\u5b66\u7684\u30ab\u30c6\u30b4\u30ea\u30fc\n(ICD, DSM)",
            ha="center", va="center", fontsize=9, fontweight="bold")

    arrow_props = dict(arrowstyle="->,head_width=0.3,head_length=0.2",
                       color="black", linewidth=2)
    ax.annotate("", xy=(5.25, 2.0), xytext=(5.25, 1.1),
                arrowprops=arrow_props)
    ax.text(5.7, 1.5, "\u30ab\u30c6\u30b4\u30ea\u30fc\u3092\n\u5b9a\u7fa9", fontsize=8, color="gray")

    ax.annotate("", xy=(2.75, 5.0), xytext=(4.0, 3.5),
                arrowprops=dict(arrowstyle="->", color="#FF6B6B",
                                linewidth=2, linestyle="--"))
    ax.text(2.2, 4.2, "\u8a3a\u65ad\u30fb\u6cbb\u7642\u3092\n\u5f62\u6210", fontsize=8,
            color="#FF6B6B", ha="center")

    ax.annotate("", xy=(7.75, 5.0), xytext=(6.5, 3.5),
                arrowprops=dict(arrowstyle="->", color="#FF6B6B",
                                linewidth=2, linestyle="--"))
    ax.text(8.2, 4.2, "\u60a3\u8005\u306e\u81ea\u5df1\u7406\u89e3\u3092\n\u5f62\u6210",
            fontsize=8, color="#FF6B6B", ha="center")

    ax.annotate("", xy=(6.5, 5.75), xytext=(4.0, 5.75),
                arrowprops=dict(arrowstyle="<->", color="gray",
                                linewidth=1.5))
    ax.text(5.25, 6.1, "\u5f93\u6765\u306e\u95a2\u4fc2", fontsize=8,
            color="gray", ha="center")

    ax.text(5.25, 6.7, "\u533b\u7642\u7248\u30b5\u30d4\u30a2\uff1d\u30a6\u30a9\u30fc\u30d5 \u30d5\u30a3\u30fc\u30c9\u30d0\u30c3\u30af\u30eb\u30fc\u30d7",
            ha="center", va="center", fontsize=14, fontweight="bold",
            color="#2C3E50")

    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight",
                facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def create_figure2_en():
    """ITS study design diagram (English)."""
    fig, ax = plt.subplots(1, 1, figsize=(10, 5))
    ax.set_xlim(2018, 2032)
    ax.set_ylim(-0.5, 4.5)
    ax.set_xlabel("Year", fontsize=12)
    ax.set_yticks([0, 1, 2, 3, 4])
    ax.set_yticklabels(["Country E", "Country D", "Country C",
                        "Country B", "Country A"], fontsize=10)
    ax.set_title("Staggered ICD-11 Adoption: Multiple Baseline Design",
                 fontsize=14, fontweight="bold")

    adoption_years = [2022, 2024, 2025, 2027, 2029]
    colors_pre = "#AED6F1"
    colors_post = "#F1948A"

    for i, yr in enumerate(adoption_years):
        ax.barh(4 - i, yr - 2018, left=2018, height=0.6,
                color=colors_pre, edgecolor="gray")
        ax.barh(4 - i, 2032 - yr, left=yr, height=0.6,
                color=colors_post, edgecolor="gray")
        ax.axvline(x=yr, ymin=(4 - i) / 5.0, ymax=(5 - i) / 5.0,
                   color="red", linewidth=2, linestyle="--")
        ax.text(yr, 4 - i + 0.35, f"ICD-11\nadopted",
                fontsize=7, ha="center", color="red")

    ax.legend(
        [mpatches.Patch(color=colors_pre), mpatches.Patch(color=colors_post)],
        ["Pre-ICD-11 (ICD-10)", "Post-ICD-11"],
        loc="upper left", fontsize=9
    )

    # Measurement arrows
    ax.annotate("", xy=(2032, -0.3), xytext=(2018, -0.3),
                arrowprops=dict(arrowstyle="<->", color="black", linewidth=1))
    ax.text(2025, -0.45, "Outcome measurement period", ha="center",
            fontsize=9, color="black")

    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight",
                facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def create_figure2_ja():
    """ITS study design diagram (Japanese)."""
    plt.rcParams["font.family"] = "Noto Sans CJK JP"
    fig, ax = plt.subplots(1, 1, figsize=(10, 5))
    ax.set_xlim(2018, 2032)
    ax.set_ylim(-0.5, 4.5)
    ax.set_xlabel("\u5e74", fontsize=12)
    ax.set_yticks([0, 1, 2, 3, 4])
    ax.set_yticklabels(["\u56fd E", "\u56fd D", "\u56fd C",
                        "\u56fd B", "\u56fd A"], fontsize=10)
    ax.set_title("ICD-11\u306e\u6bb5\u968e\u7684\u63a1\u7528\uff1a\u591a\u91cd\u30d9\u30fc\u30b9\u30e9\u30a4\u30f3\u30c7\u30b6\u30a4\u30f3",
                 fontsize=14, fontweight="bold")

    adoption_years = [2022, 2024, 2025, 2027, 2029]
    colors_pre = "#AED6F1"
    colors_post = "#F1948A"

    for i, yr in enumerate(adoption_years):
        ax.barh(4 - i, yr - 2018, left=2018, height=0.6,
                color=colors_pre, edgecolor="gray")
        ax.barh(4 - i, 2032 - yr, left=yr, height=0.6,
                color=colors_post, edgecolor="gray")
        ax.axvline(x=yr, ymin=(4 - i) / 5.0, ymax=(5 - i) / 5.0,
                   color="red", linewidth=2, linestyle="--")
        ax.text(yr, 4 - i + 0.35, "ICD-11\n\u63a1\u7528",
                fontsize=7, ha="center", color="red")

    ax.legend(
        [mpatches.Patch(color=colors_pre), mpatches.Patch(color=colors_post)],
        ["ICD-11\u524d (ICD-10)", "ICD-11\u5f8c"],
        loc="upper left", fontsize=9
    )

    ax.annotate("", xy=(2032, -0.3), xytext=(2018, -0.3),
                arrowprops=dict(arrowstyle="<->", color="black", linewidth=1))
    ax.text(2025, -0.45, "\u30a2\u30a6\u30c8\u30ab\u30e0\u6e2c\u5b9a\u671f\u9593", ha="center",
            fontsize=9, color="black")

    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight",
                facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def add_slide(prs, title_text, caption_text, image_buf):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.8)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Image
    slide.shapes.add_picture(
        image_buf, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.0)
    )

    # Caption
    txBox2 = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.4), Inches(12.333), Inches(1.0)
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = caption_text
    p2.font.size = Pt(12)
    p2.font.italic = True
    p2.alignment = PP_ALIGN.CENTER


def add_table_slide(prs, title_text, headers, rows, caption_text):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    txBox = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.15), Inches(12.733), Inches(0.6)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(20)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.3), Inches(0.85),
        Inches(12.733), Inches(5.5)
    )
    tbl = tbl_shape.table

    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for par in cell.text_frame.paragraphs:
            par.font.size = Pt(11)
            par.font.bold = True
            par.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)
        for par in cell.text_frame.paragraphs:
            par.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for r, row_data in enumerate(rows, start=1):
        for c, val in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = val
            for par in cell.text_frame.paragraphs:
                par.font.size = Pt(9)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Caption
    txBox2 = slide.shapes.add_textbox(
        Inches(0.3), Inches(6.5), Inches(12.733), Inches(0.8)
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = caption_text
    p2.font.size = Pt(10)
    p2.font.italic = True
    p2.alignment = PP_ALIGN.CENTER


def main():
    # ====================== ENGLISH ======================
    prs_en = Presentation()
    prs_en.slide_width = SLIDE_W
    prs_en.slide_height = SLIDE_H

    # Figure 1 EN
    fig1_en = create_figure1_en()
    add_slide(
        prs_en,
        "Figure 1. Medical Sapir\u2013Whorf Feedback Loop",
        "Nosological categories operate at the sickness level, shaping institutional "
        "responses. These changes feed back to alter disease recognition and illness "
        "experience, creating a recursive loop.",
        fig1_en
    )

    # Table 1 EN
    add_table_slide(
        prs_en,
        "Table 1. Clinical Responses to Work-Related CV Events: Japan vs. Countries Without Karoshi",
        ["Domain", "Japan (karoshi concept present)",
         "Countries without equivalent concept"],
        [
            ["Diagnostic pathway",
             "Work history systematically assessed; occupational medicine consultation",
             "Standard CVA/MI workup; working conditions rarely documented"],
            ["Coding",
             "Specific codes link CV event to overwork",
             "Standard ICD CV codes; no work linkage"],
            ["Compensation",
             "Dedicated workers\u2019 compensation with legal criteria (>80 h/week)",
             "General disability insurance; no work-specific CV pathway"],
            ["Epidemiological tracking",
             "National karoshi statistics since 1987",
             "No systematic tracking of work-related CV mortality"],
            ["Prevention infrastructure",
             "Mandatory overtime limits, annual health checks, karoshi hotline",
             "General OH regulations; no targeted CV prevention"],
            ["Research funding",
             "Dedicated karoshi research; large epidemiological datasets",
             "Work-related CVD as subtopic of occupational health"],
        ],
        "Table 1. Comparison illustrating the infrastructure creation mechanism."
    )

    # Figure 2 EN
    fig2_en = create_figure2_en()
    add_slide(
        prs_en,
        "Figure 2. Proposed Interrupted Time-Series Study Design",
        "Staggered ICD-11 adoption across countries enables a multiple baseline "
        "design, strengthening causal inference about classification changes on "
        "clinical practice.",
        fig2_en
    )

    # Table 2 EN
    add_table_slide(
        prs_en,
        "Table 2. Key ICD-11 Structural Changes for Sapir\u2013Whorf Analysis",
        ["ICD-11 change", "What can be tested",
         "Expected Sapir\u2013Whorf effect"],
        [
            ["Chronic pain chapter (MG30)",
             "Referral patterns, specialist utilization, treatment modalities",
             "Pain recognized as disease; increased specialist referrals"],
            ["Burnout (QD85)",
             "International diffusion of karoshi concepts; sick-leave patterns",
             "Formal classification catalyzes karoshi-like infrastructure"],
            ["Gaming disorder (6C51)",
             "Incidence trends; treatment-seeking behavior",
             "Classification creates clinical entity (looping effect)"],
            ["Gender incongruence reclassified",
             "Stigma, treatment access, self-identification",
             "Reclassification alters illness experience without changing biology"],
        ],
        "Table 2. ICD-11 changes providing natural experimental conditions."
    )

    en_path = os.path.join(OUT_DIR, "figures_tables_en.pptx")
    prs_en.save(en_path)
    print(f"Saved: {en_path}")

    # ====================== JAPANESE ======================
    prs_ja = Presentation()
    prs_ja.slide_width = SLIDE_W
    prs_ja.slide_height = SLIDE_H

    fig1_ja = create_figure1_ja()
    add_slide(
        prs_ja,
        "\u56f31. \u533b\u7642\u7248\u30b5\u30d4\u30a2\uff1d\u30a6\u30a9\u30fc\u30d5 \u30d5\u30a3\u30fc\u30c9\u30d0\u30c3\u30af\u30eb\u30fc\u30d7",
        "\u75be\u75c5\u5206\u985e\u5b66\u7684\u30ab\u30c6\u30b4\u30ea\u30fc\u306f\u75c5\u4eba\u5f79\u5272\u30ec\u30d9\u30eb\u3067\u4f5c\u52d5\u3057\u3001\u5236\u5ea6\u7684\u5fdc\u7b54\u3092\u5f62\u6210\u3059\u308b\u3002\u3053\u306e\u5909\u5316\u306f\u75be\u60a3\u8a8d\u8b58\u304a\u3088\u3073"
        "\u75c5\u3044\u4f53\u9a13\u306b\u30d5\u30a3\u30fc\u30c9\u30d0\u30c3\u30af\u3057\u3001\u518d\u5e30\u7684\u30eb\u30fc\u30d7\u3092\u5275\u51fa\u3059\u308b\u3002",
        fig1_ja
    )

    add_table_slide(
        prs_ja,
        "\u88681. \u52b4\u50cd\u95a2\u9023\u5fc3\u8840\u7ba1\u30a4\u30d9\u30f3\u30c8\u3078\u306e\u81e8\u5e8a\u7684\u5fdc\u7b54\uff1a\u65e5\u672c vs. \u904e\u52b4\u6b7b\u6982\u5ff5\u306a\u3057\u306e\u56fd",
        ["\u9818\u57df", "\u65e5\u672c\uff08\u904e\u52b4\u6b7b\u6982\u5ff5\u3042\u308a\uff09",
         "\u5bfe\u5fdc\u6982\u5ff5\u306a\u3057\u306e\u56fd"],
        [
            ["\u8a3a\u65ad\u30d1\u30b9\u30a6\u30a7\u30a4",
             "\u52b4\u50cd\u6b74\u304c\u4f53\u7cfb\u7684\u306b\u8a55\u4fa1\uff1b\u7523\u696d\u533b\u30b3\u30f3\u30b5\u30eb\u30c6\u30fc\u30b7\u30e7\u30f3\u304c\u767a\u52d5",
             "\u6a19\u6e96\u7684CVA/MI\u7cbe\u67fb\uff1b\u52b4\u50cd\u6761\u4ef6\u306f\u7a00\u306b\u8a18\u9332"],
            ["\u30b3\u30fc\u30c7\u30a3\u30f3\u30b0",
             "\u5fc3\u8840\u7ba1\u30a4\u30d9\u30f3\u30c8\u3092\u904e\u91cd\u52b4\u50cd\u3068\u7d10\u4ed8\u3051\u308b\u7279\u5b9a\u30b3\u30fc\u30c9",
             "\u6a19\u6e96ICD\u5fc3\u8840\u7ba1\u30b3\u30fc\u30c9\uff1b\u52b4\u50cd\u3068\u306e\u7d10\u4ed8\u304d\u306a\u3057"],
            ["\u88dc\u511f",
             "\u6cd5\u7684\u57fa\u6e96\u3092\u4f34\u3046\u5c02\u7528\u52b4\u707d\u88dc\u511f\uff08\u4f8b\uff1a\u670880\u6642\u9593\u8d85\u306e\u6b8b\u696d\uff09",
             "\u4e00\u822c\u7684\u969c\u5bb3\u4fdd\u967a\uff1b\u696d\u52d9\u7279\u7570\u7684\u30d1\u30b9\u30a6\u30a7\u30a4\u306a\u3057"],
            ["\u75ab\u5b66\u7684\u8ffd\u8de1",
             "1987\u5e74\u4ee5\u964d\u306e\u904e\u52b4\u6b7b\u75c7\u4f8b\u306e\u5168\u56fd\u7d71\u8a08",
             "\u52b4\u50cd\u95a2\u9023\u5fc3\u8840\u7ba1\u6b7b\u4ea1\u306e\u4f53\u7cfb\u7684\u8ffd\u8de1\u306a\u3057"],
            ["\u4e88\u9632\u30a4\u30f3\u30d5\u30e9",
             "\u6642\u9593\u5916\u52b4\u50cd\u4e0a\u9650\u898f\u5236\u3001\u5e74\u6b21\u5065\u8a3a\u3001\u904e\u52b4\u6b7b110\u756a",
             "\u4e00\u822c\u7684\u7523\u696d\u885b\u751f\u898f\u5236\uff1b\u6a19\u7684\u578b\u5fc3\u8840\u7ba1\u4e88\u9632\u306a\u3057"],
            ["\u7814\u7a76\u8cc7\u91d1",
             "\u904e\u52b4\u6b7b\u5c02\u9580\u306e\u7814\u7a76\uff1b\u5927\u898f\u6a21\u75ab\u5b66\u30c7\u30fc\u30bf",
             "\u52b4\u50cd\u95a2\u9023CVD\u306f\u7523\u696d\u885b\u751f\u306e\u30b5\u30d6\u30c8\u30d4\u30c3\u30af"],
        ],
        "\u88681. \u30a4\u30f3\u30d5\u30e9\u30b9\u30c8\u30e9\u30af\u30c1\u30e3\u30fc\u5275\u51fa\u30e1\u30ab\u30cb\u30ba\u30e0\u3092\u4f8b\u793a\u3059\u308b\u6bd4\u8f03\u3002"
    )

    fig2_ja = create_figure2_ja()
    add_slide(
        prs_ja,
        "\u56f32. \u63d0\u6848\u3059\u308b\u5206\u65ad\u6642\u7cfb\u5217\u7814\u7a76\u30c7\u30b6\u30a4\u30f3",
        "\u56fd\u5bb6\u9593\u306eICD-11\u306e\u6bb5\u968e\u7684\u63a1\u7528\u304c\u591a\u91cd\u30d9\u30fc\u30b9\u30e9\u30a4\u30f3\u30c7\u30b6\u30a4\u30f3\u3092\u53ef\u80fd\u306b\u3057\u3001"
        "\u5206\u985e\u5909\u66f4\u306e\u81e8\u5e8a\u5b9f\u8df5\u3078\u306e\u52b9\u679c\u306b\u95a2\u3059\u308b\u56e0\u679c\u63a8\u8ad6\u3092\u5f37\u5316\u3059\u308b\u3002",
        fig2_ja
    )

    add_table_slide(
        prs_ja,
        "\u88682. \u30b5\u30d4\u30a2\uff1d\u30a6\u30a9\u30fc\u30d5\u5206\u6790\u306b\u9069\u3059\u308bICD-11\u306e\u4e3b\u8981\u69cb\u9020\u5909\u5316",
        ["ICD-11\u306e\u5909\u66f4", "\u691c\u8a3c\u53ef\u80fd\u306a\u4e8b\u9805",
         "\u4e88\u6e2c\u3055\u308c\u308b\u30b5\u30d4\u30a2\uff1d\u30a6\u30a9\u30fc\u30d5\u52b9\u679c"],
        [
            ["\u6162\u6027\u75bc\u75db\u72ec\u7acb\u7ae0\uff08MG30\uff09",
             "\u7d39\u4ecb\u30d1\u30bf\u30fc\u30f3\u30fb\u5c02\u9580\u533b\u5229\u7528\u30fb\u6cbb\u7642\u6cd5\u306e\u5909\u5316",
             "\u75bc\u75db\u304c\u75be\u60a3\u5b9f\u4f53\u3068\u3057\u3066\u8a8d\u8b58\uff1b\u5c02\u9580\u533b\u7d39\u4ecb\u306e\u5897\u52a0"],
            ["\u30d0\u30fc\u30f3\u30a2\u30a6\u30c8\uff08QD85\uff09",
             "\u904e\u52b4\u6b7b\u985e\u4f3c\u6982\u5ff5\u306e\u56fd\u969b\u7684\u62e1\u6563\uff1b\u75c5\u6c17\u4f11\u6687\u30d1\u30bf\u30fc\u30f3\u306e\u5909\u5316",
             "\u516c\u5f0f\u5206\u985e\u304c\u904e\u52b4\u6b7b\u985e\u4f3c\u306e\u30a4\u30f3\u30d5\u30e9\u3092\u89e6\u5a92"],
            ["\u30b2\u30fc\u30e0\u969c\u5bb3\uff086C51\uff09",
             "\u767a\u751f\u7387\u306e\u63a8\u79fb\uff1b\u53d7\u7642\u884c\u52d5",
             "\u5206\u985e\u304c\u81e8\u5e8a\u7684\u5b9f\u4f53\u3092\u5275\u51fa\uff08\u30eb\u30fc\u30d7\u52b9\u679c\uff09"],
            ["\u6027\u5225\u4e0d\u5408\u306e\u518d\u5206\u985e",
             "\u30b9\u30c6\u30a3\u30b0\u30de\u30fb\u6cbb\u7642\u30a2\u30af\u30bb\u30b9\u30fb\u81ea\u5df1\u540c\u5b9a\u306e\u5909\u5316",
             "\u751f\u7269\u5b66\u3092\u5909\u3048\u305a\u306b\u75c5\u3044\u4f53\u9a13\u3092\u5909\u5bb9"],
        ],
        "\u88682. \u81ea\u7136\u5b9f\u9a13\u6761\u4ef6\u3092\u63d0\u4f9b\u3059\u308bICD-11\u306e\u5909\u66f4\u3002"
    )

    ja_path = os.path.join(OUT_DIR, "figures_tables_ja.pptx")
    prs_ja.save(ja_path)
    print(f"Saved: {ja_path}")


if __name__ == "__main__":
    main()
